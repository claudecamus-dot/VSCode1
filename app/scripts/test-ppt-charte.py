"""Tests fonctionnels de charte graphique du PPT de restitution (US6.4).

Complementaires a test-export-ppt.py (structure + geometrie) : ici on verifie
la QUALITE visuelle du rendu — ce qu'une revue humaine ("pptx-verify") ferait
a l'oeil, mais rejoue en assertions sur le contenu reel du .pptx genere avec
le VRAI template OCTO (pas de fixture de theme factice) :

  - police de marque appliquee partout sur le texte dessine (pas de police
    etrangere qui trahirait un repli Arial) ;
  - tailles de police dans des bornes de lisibilite (ni ecrasees, ni demesurees) ;
  - couleurs utilisees dans la palette approuvee (charte + palette piliers +
    exceptions documentees) — detecte une couleur ad hoc introduite par erreur ;
  - contraste texte/fond conforme WCAG AA (4.5:1 texte normal, 3:1 grand texte
    gras >= 14pt) pour les couleurs de texte reellement utilisees ;
  - alignement "tableau" des lignes (barres par pilier, evolution) : colonnes
    alignees, lignes non chevauchantes, espacement regulier.

Usage : python test-ppt-charte.py
"""
import importlib.util
import os
import sys
import tempfile

from pptx.enum.dml import MSO_COLOR_TYPE, MSO_FILL_TYPE
from pptx.enum.shapes import MSO_SHAPE_TYPE

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))


def _charger(nom_fichier, nom_module):
    """Les generateurs ont un nom de fichier non importable (tirets)."""
    chemin = os.path.join(os.path.dirname(os.path.abspath(__file__)), nom_fichier)
    spec = importlib.util.spec_from_file_location(nom_module, chemin)
    mod = importlib.util.module_from_spec(spec)
    spec.loader.exec_module(mod)
    return mod


gen = _charger("export-restitution-ppt.py", "gen_charte")
fixtures = _charger("test-export-ppt.py", "fixtures_charte")
D = gen.D

echecs = 0


def check(cond, msg):
    global echecs
    if cond:
        print(f"  ok   {msg}")
    else:
        echecs += 1
        print(f"  FAIL {msg}")


# ----------------------------------------------------------------------------
# Construction d'un deck realiste (vrai template OCTO, donnees normales — pas
# le cas de stress geometrique de test-export-ppt.py) pour l'inspecter.
# ----------------------------------------------------------------------------
def _construire_deck():
    tmp = tempfile.mkdtemp(prefix="test-ppt-charte-")
    bloc = fixtures.bloc_equipe("DSI Paiements", [2.0, 1.8, 1.6, 1.9], avec_comp=True)
    data = {
        "couverture": {"titre": "Restitution — Maturité agile/produit",
                       "sousTitre": "DSI Paiements", "date": "08/07/2026"},
        "blocs": [bloc],
    }
    out = os.path.join(tmp, "deck.pptx")
    prs, problemes = gen.construire(data, gen.TEMPLATE, out)
    return prs, problemes


def _slide_texte_boxes(slide):
    """Zones de texte reellement DESSINEES par notre code (D.add_text /
    add_textbox direct) — exclut les placeholders du template (couverture),
    qui heritent de la charte du modele et ne sont pas de notre ressort ici."""
    return [s for s in slide.shapes
            if s.shape_type == MSO_SHAPE_TYPE.TEXT_BOX and s.has_text_frame]


def _runs_non_vides(shp):
    for para in shp.text_frame.paragraphs:
        for run in para.runs:
            if run.text.strip():
                yield run


# ----------------------------------------------------------------------------
# 1. Police de marque
# ----------------------------------------------------------------------------
def test_police_marque(prs):
    marque = D.POLICE
    check(bool(marque), f"police de marque detectee dans le template (detecte : {marque!r})")
    if marque:
        repli_generiques = {"arial", "calibri", "times new roman", "segoe ui"}
        check(marque.lower() not in repli_generiques,
              f"la police detectee n'est pas un repli generique (detecte : {marque})")

    autres_polices = set()
    n_runs = 0
    for slide in prs.slides:
        for shp in _slide_texte_boxes(slide):
            for run in _runs_non_vides(shp):
                n_runs += 1
                if run.font.name and run.font.name != marque:
                    autres_polices.add(run.font.name)
    check(n_runs > 0, f"au moins un run de texte dessine trouve ({n_runs})")
    check(not autres_polices,
          "toutes les zones de texte dessinees utilisent la police de marque"
          if not autres_polices else
          f"police(s) etrangere(s) trouvee(s) sur le texte dessine : {sorted(autres_polices)}")


# ----------------------------------------------------------------------------
# 2. Tailles de police dans des bornes de lisibilite
# ----------------------------------------------------------------------------
# Plancher legerement sous TAILLE_MIN_CARTE (repli extreme de l'auto-ajustement
# des cartes) pour ne pas fausser-positiver sur ce cas deja couvert par
# test-export-ppt.py ; plafond = plus grande taille de l'echelle (kpi).
PLANCHER_PT = gen.TAILLE_MIN_CARTE - 0.5
PLAFOND_PT = D.TYPE["kpi"]


def test_tailles_police(prs):
    tailles = []
    for slide in prs.slides:
        for shp in _slide_texte_boxes(slide):
            for run in _runs_non_vides(shp):
                if run.font.size is not None:
                    tailles.append(run.font.size.pt)
    check(bool(tailles), f"au moins une taille de police mesuree ({len(tailles)} runs)")
    hors_bornes = sorted({t for t in tailles if t < PLANCHER_PT or t > PLAFOND_PT})
    check(not hors_bornes,
          f"toutes les tailles dans [{PLANCHER_PT}, {PLAFOND_PT}] pt"
          if not hors_bornes else
          f"taille(s) hors bornes [{PLANCHER_PT}, {PLAFOND_PT}] pt : {hors_bornes}")


# ----------------------------------------------------------------------------
# 3. Couleurs dans la palette approuvee
# ----------------------------------------------------------------------------
def _palette_autorisee():
    autorises = {D.INK, D.MUTED, D.LINE, D.TRACK, D.CYAN, gen.ACCENT,
                 gen.FOND_PANNEAU, D.OK, D.WARN, D.GOLD, "#ffffff"}
    autorises |= set(D.PALETTE)
    # Teinte ad hoc assumee : rouge adouci pour la barre "score le plus faible"
    # (distinct de D.WARN, utilise uniquement sur la valeur) — voir
    # export-restitution-ppt.py:rendu_faible. Documentee ici plutot que
    # silencieusement ignoree : c'est une exception connue, pas la charte.
    autorises.add("#cf7b74")
    return {c.upper() for c in autorises}


def _couleur_fill(shp):
    try:
        if shp.fill.type == MSO_FILL_TYPE.SOLID:
            return "#" + str(shp.fill.fore_color.rgb)
    except (TypeError, AttributeError, ValueError):
        pass
    return None


def _couleur_ligne(shp):
    try:
        if shp.line.fill.type == MSO_FILL_TYPE.SOLID:
            return "#" + str(shp.line.color.rgb)
    except (TypeError, AttributeError, ValueError):
        pass
    return None


def _couleur_font(run):
    try:
        if run.font.color and run.font.color.type == MSO_COLOR_TYPE.RGB:
            return "#" + str(run.font.color.rgb)
    except (TypeError, AttributeError, ValueError):
        pass
    return None


def test_couleurs_autorisees(prs):
    autorises = _palette_autorisee()
    trouvees = set()
    for slide in prs.slides:
        for shp in slide.shapes:
            for getter in (_couleur_fill, _couleur_ligne):
                c = getter(shp)
                if c:
                    trouvees.add(c.upper())
            if shp.has_text_frame:
                for para in shp.text_frame.paragraphs:
                    for run in para.runs:
                        c = _couleur_font(run)
                        if c:
                            trouvees.add(c.upper())
    inconnues = sorted(trouvees - autorises)
    check(not inconnues,
          f"toutes les couleurs utilisees sont dans la palette approuvee ({len(trouvees)} distinctes)"
          if not inconnues else
          f"couleur(s) hors palette approuvee : {inconnues}")


# ----------------------------------------------------------------------------
# 4. Contraste texte/fond (WCAG AA) pour les couleurs de texte reellement
#    utilisees dans le deck, sur leur fond reel (blanc de la slide, ou panneau
#    FOND_PANNEAU pour le texte en encart).
# ----------------------------------------------------------------------------
def _luminance(hexcolor):
    hexcolor = hexcolor.lstrip("#")
    vals = [int(hexcolor[i:i + 2], 16) / 255 for i in (0, 2, 4)]

    def lin(c):
        return c / 12.92 if c <= 0.03928 else ((c + 0.055) / 1.055) ** 2.4

    r, g, b = (lin(c) for c in vals)
    return 0.2126 * r + 0.7152 * g + 0.0722 * b


def _contraste(hex1, hex2):
    l1, l2 = sorted((_luminance(hex1), _luminance(hex2)), reverse=True)
    return (l1 + 0.05) / (l2 + 0.05)


def _seuil_wcag(taille_pt, gras):
    """WCAG AA : 3:1 pour un 'grand texte' (>=18pt, ou >=14pt gras), 4.5:1 sinon."""
    grand = taille_pt >= 18 or (gras and taille_pt >= 14)
    return 3.0 if grand else 4.5


# Couples (libelle, couleur, taille pt, gras, fond) reellement utilises par le
# generateur pour du texte colore (hors INK/MUTED, deja tres contrastes) —
# lus sur les VRAIES valeurs post-theme (D.*, gen.ACCENT/FOND_PANNEAU), pas
# des constantes recopiees a la main. Fonds : blanc (texte pose a meme la
# slide, ex. libelles d'axe du radar) ou FOND_PANNEAU (texte en encart).
def _couples_a_verifier():
    blanc = "#ffffff"
    panneau = gen.FOND_PANNEAU
    couples = []
    # Libelles d'axe du radar : depuis 2026-07-21 en D.INK (fonce neutre, deja tres
    # contraste — non teste ici, comme INK/MUTED). La couleur du pilier est passee du
    # TEXTE a la PASTILLE (objet graphique) : le gold #b8860b echouait en texte
    # (3.25:1 < AA 4.5:1) mais passe comme pastille (seuil graphique 3:1). Le contraste
    # des libelles ET des 6 pastilles pilier, sur les DEUX surfaces (web + PPT, depuis
    # la meme palette), est verifie par scripts/test-contraste-radar.js.
    # Valeurs colorees (ecart-type, scores) : h3 (14pt) gras, sur blanc (carte).
    for nom, coul in (("OK", D.OK), ("WARN", D.WARN), ("GOLD", D.GOLD)):
        couples.append((f"valeur carte {nom} ({coul}) sur blanc",
                        coul, D.TYPE["h3"], True, blanc))
    # Texte MUTED (libelles secondaires) : small (10.5pt), sur blanc et sur panneau.
    couples.append((f"texte secondaire MUTED ({D.MUTED}) sur blanc", D.MUTED, D.TYPE["small"], False, blanc))
    couples.append((f"texte secondaire MUTED ({D.MUTED}) sur panneau", D.MUTED, D.TYPE["small"], False, panneau))
    return couples


def test_contraste_wcag(prs):
    for libelle, coul, taille, gras, fond in _couples_a_verifier():
        seuil = _seuil_wcag(taille, gras)
        ratio = _contraste(coul, fond)
        check(ratio >= seuil, f"{libelle} : contraste {ratio:.2f}:1 (seuil WCAG AA {seuil}:1)")


# ----------------------------------------------------------------------------
# 5. Rendu "tableau" : colonnes alignees, lignes non chevauchantes — verifie
#    par la geometrie des shapes (regroupement par colonne `left`, pas par
#    hauteur : une ligne peut legitimement etre plus haute qu'une autre quand
#    son libelle se replie sur plusieurs lignes, cf. slide_radar). Robuste a
#    un refactor de layout tant que les colonnes restent alignees verticalement.
# ----------------------------------------------------------------------------
def _colonnes_candidates(boxes, n_attendu, tol_in=0.003):
    """Groupe `boxes` par position horizontale (`left`, arrondi) ; ne garde
    que les groupes de taille EXACTEMENT `n_attendu` (une ligne par element
    attendu) — ecarte tout le reste (titres, legendes, echelle...)."""
    par_gauche = {}
    for b in boxes:
        cle = round(b.left.inches / tol_in) * tol_in
        par_gauche.setdefault(cle, []).append(b)
    return [sorted(v, key=lambda b: b.top) for v in par_gauche.values() if len(v) == n_attendu]


def _assert_colonnes_alignees(nom_slide, colonnes, n_lignes):
    """`colonnes` = 2+ listes de `n_lignes` zones de texte (une par colonne,
    triees par `top`). Verifie : les colonnes s'alignent ligne a ligne (meme
    `top` a l'index i pour toutes les colonnes) et les lignes se succedent
    sans chevauchement (chaque ligne se termine avant que la suivante ne
    commence, en tenant compte de sa PROPRE hauteur — qui peut varier)."""
    check(len(colonnes) >= 2,
          f"{nom_slide} : au moins 2 colonnes de {n_lignes} ligne(s) detectees (trouve {len(colonnes)})")
    if len(colonnes) < 2:
        return
    ref = [round(b.top.inches, 3) for b in colonnes[0]]
    for i, col in enumerate(colonnes[1:], start=1):
        tops = [round(b.top.inches, 3) for b in col]
        check(tops == ref, f"{nom_slide} : colonne {i} alignee ligne a ligne sur la colonne 0")
    hauteurs = [round(b.height.inches, 3) for b in colonnes[0]]
    # Tolerance de 0.005in (~0.13mm) : sous ce seuil, c'est du bruit
    # d'arrondi flottant (ex. l'echelle proportionnelle appliquee quand le
    # total deborde, cf. slide_radar), pas un chevauchement visible.
    TOL = 0.005
    chevauchements = [(t2 < t1 + h1 - TOL) for t1, h1, t2 in zip(ref, hauteurs, ref[1:])]
    check(not any(chevauchements),
          f"{nom_slide} : lignes non chevauchantes (tops {ref}, hauteurs {hauteurs})")
    check(all(b - a > 0 for a, b in zip(ref, ref[1:])),
          f"{nom_slide} : lignes dans l'ordre (tops strictement croissants)")


def _trouver_slide(prs, *marqueurs):
    """1ere slide dont un shape (textbox OU placeholder) contient tous les
    `marqueurs` (recherche large : le titre vit dans un placeholder de layout,
    pas un textbox dessine par notre code — contrairement au reste du texte)."""
    for s in prs.slides:
        textes = " ".join(
            r.text for shp in s.shapes if shp.has_text_frame
            for para in shp.text_frame.paragraphs for r in para.runs if r.text.strip())
        if all(m in textes for m in marqueurs):
            return s
    return None


def test_tableau_barres_piliers(prs):
    slide = _trouver_slide(prs, "Vue d'ensemble", "MATURITÉ PAR PILIER")
    check(slide is not None, "slide 'Vue d'ensemble' trouvee")
    if slide is None:
        return
    n_piliers = len(fixtures.piliers([0, 0, 0, 0]))
    colonnes = _colonnes_candidates(_slide_texte_boxes(slide), n_piliers)
    # 2 colonnes attendues (libelle pilier, valeur) parmi les groupes candidats
    # (d'autres largeurs de n_piliers boites peuvent apparaitre par coincidence,
    # ex. la reglette d'echelle a 4 graduations -> on prend les 2 premieres
    # colonnes triees par position horizontale, qui sont les plus a gauche/
    # constantes sur toute la hauteur de la bande de barres).
    colonnes = sorted(colonnes, key=lambda col: col[0].left.inches)[:2]
    _assert_colonnes_alignees("Vue d'ensemble (barres pilier)", colonnes, n_piliers)


def test_tableau_evolution(prs):
    slide = _trouver_slide(prs, "ÉVOLUTION VS")
    check(slide is not None, "slide radar avec bloc 'evolution' trouvee")
    if slide is None:
        return
    # La legende du radar (a gauche) partage parfois le meme nombre de lignes
    # que le nombre de piliers -> on scope la recherche au panneau evolution,
    # repere par le marqueur textuel "ÉVOLUTION VS..." (son `left` = debut du
    # panneau) plutot que par une coordonnee codee en dur.
    marqueur = next((shp for shp in _slide_texte_boxes(slide)
                     for r in _runs_non_vides(shp) if "ÉVOLUTION VS" in r.text.upper()), None)
    check(marqueur is not None, "marqueur 'ÉVOLUTION VS' localise sur la slide")
    if marqueur is None:
        return
    px = marqueur.left.inches - 0.01   # petite tolerance
    boxes = [b for b in _slide_texte_boxes(slide) if b.left.inches >= px]
    n_piliers = len(fixtures.piliers([0, 0, 0, 0]))
    colonnes = _colonnes_candidates(boxes, n_piliers)
    check(len(colonnes) >= 2,
          f"Radar (evolution par pilier) : au moins 2 colonnes de {n_piliers} ligne(s) candidates "
          f"(trouve {len(colonnes)})")
    if len(colonnes) < 2:
        return
    colonnes = sorted(colonnes, key=lambda col: col[0].left.inches)
    _assert_colonnes_alignees("Radar (evolution par pilier)", colonnes, n_piliers)


def main():
    print("Construction du deck (vrai template OCTO) :")
    prs, problemes = _construire_deck()
    check(not problemes, f"geometrie OK (prealable a ces tests) — {len(problemes)} probleme(s)")

    print("\nPolice de marque :")
    test_police_marque(prs)

    print("\nTailles de police :")
    test_tailles_police(prs)

    print("\nCouleurs autorisees :")
    test_couleurs_autorisees(prs)

    print("\nContraste texte/fond (WCAG AA) :")
    test_contraste_wcag(prs)

    print("\nRendu tableau (alignement des lignes) :")
    test_tableau_barres_piliers(prs)
    test_tableau_evolution(prs)

    print("\nTOUS LES TESTS PASSENT" if echecs == 0 else f"\n{echecs} TEST(S) EN ECHEC")
    sys.exit(0 if echecs == 0 else 1)


if __name__ == "__main__":
    main()
