"""Genere un support de restitution PPT (template OCTO) — style INFOGRAPHIE (US6.4).

Usage : python export-restitution-ppt.py <donnees.json> <sortie.pptx> [modele.pptx]

Le template du modele est, par ordre de priorite : l'argument [modele.pptx], la
variable d'env TEMPLATE_PPTX, sinon le template OCTO par defaut. La couverture et le
layout "titre seul" sont reperes par NOM (repli sur indices), et la couleur d'accent
de marque (jauge, callout) est derivee du theme du modele — voir construire().

Structure : 1 couverture (layout OCTO inchange), puis pour chaque "bloc" (une
equipe, ou un departement consolidant >= 2 equipes) 4 slides :
  1. Vue d'ensemble  — jauge "moyenne globale" + barres de maturite par pilier
                       (couleurs du radar) + chips point fort / a renforcer.
  2. Radar           — radar par objectif (image SVG facon web) + commentaire de
                       restitution en encart + evolution par pilier.
  3. Points forts    — cartes (vert) : scores les plus hauts et meilleurs
                       accords (dispersion la plus faible), pendant positif de
                       la slide suivante.
  4. Points attention— cartes : plus forts desaccords (dispersion) et scores
                       les plus faibles.

On garde le chrome OCTO (logo, pied de page, n° de slide) via le layout "Titre
seul" ; l'infographie est dessinee dans la zone de contenu. Un controle
geometrique (pptx_deck.verifier_geometrie) garantit qu'aucune forme ne deborde.
Les cartes de points forts/attention adaptent leur police a la longueur des
phrases (pptx_deck.ajuster_police) : un cap fixe a 2 lignes faisait deborder
les questions longues du referentiel — voir _cartes_colonne.

Le serveur Node calcule tout et passe le JSON ; ce script ne touche pas a la base.
"""
import sys
import os
import json
import math
import re
from pptx import Presentation
from pptx.util import Inches
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
import pptx_deck as D

# Echelle typographique du deck. Reduction globale ×0.9 ABANDONNEE (revue design
# 2026-07-22, arbitrage utilisateur) : elle rapetissait le texte lisible (contexte,
# legende, libelles radar) SANS agrandir les cartes (deja bloquees au plancher 6pt a
# toute echelle) — donc contre-productive pour la lisibilite. Defaut = 1.0 (inchange) ;
# le knob env reste, utile pour re-comparer une magnitude sans toucher au code. pptx_deck
# n'a que ce deck pour consommateur, donc muter TYPE ici n'affecte aucun autre livrable.
_ECHELLE_POLICE = float(os.environ.get("DECK_FONT_SCALE", "1.0"))
if _ECHELLE_POLICE != 1.0:
    D.TYPE = {k: round(v * _ECHELLE_POLICE, 2) for k, v in D.TYPE.items()}
# Taille du titre de slide : le placeholder du template le fige (~28pt), non regi par TYPE.
# Abaissee a 24pt — le titre etait le seul element vraiment surdimensionne ; on garde
# CETTE reduction (arbitrage utilisateur) meme apres l'abandon de la reduction globale.
TITRE_SLIDE_PT = 24.0

TEMPLATE = os.path.join(os.path.dirname(__file__), "..", "..", "template ppt", "template.pptx")
# Layouts repérés par NOM (robuste si on fournit un autre template dont l'ordre
# des layouts diffère), avec repli sur l'indice du template OCTO d'origine.
LAYOUT_COUVERTURE = 8   # repli : "40 - Couverture [1]"
LAYOUT_TITRE_SEUL = 5   # repli : "04 - Titre seul" : garde logo/footer, titre OCTO
COUV_PATTERNS = ("couverture", "cover", "title slide")
TITRE_PATTERNS = ("titre seul", "title only", "title-only")

# Zone de contenu sous le titre (slide 10 x 5.625 in).
CONTENU_TOP = 1.15
CONTENU_BOTTOM = 5.45
MARGE_X = 0.55
RADIUS = 0.08          # rayon des coins arrondis — coherent sur tout le deck
FOND_PANNEAU = "#f7f8fb"   # gris tres clair des encarts
# Bord droit "sur" : tout contenu pleine largeur du bas (bandeau, cartes) s'arrete
# ici pour degager le badge n° de slide du template OCTO (coin bas-droit), qu'il
# recouvrirait sinon. Meme limite que le cluster numerique de la slide radar.
BORD_DROIT = 9.15
# Couleur d'accent de marque (jauge globale, barre d'accent du callout). Par defaut
# le bleu de la palette ; construire() la remplace par la couleur primaire du theme
# du template fourni (charte du modele). Les couleurs de pilier, elles, restent
# alignees sur le radar et ne sont PAS derivees du theme.
ACCENT = D.PALETTE[0]
# Accent secondaire de marque (filets de section, barre du callout) : le cyan de
# la charte. Sourcé du thème du template par construire() (D.appliquer_theme).
CYAN = D.CYAN
# Hauteur reelle d'une ligne de question (small 10.5pt) dans une carte : ~0.17in
# de texte + un peu de marge. L'ancienne valeur (0.235) sur-estimait et creusait
# un vide entre la question et le couple contexte/barre.
LH_QUESTION = 0.195


def fmt(x):
    return "—" if x is None else f"{x:.1f}"


# Les libelles de piliers viennent de l'Excel en "Title Case" maladroit
# ("Culture De L'Entreprise Agile", "Agilite A L'Echelle") : peu lisibles. On
# les remet en casse de phrase FR — seul le premier mot (et les acronymes connus)
# prend une majuscule — et on restaure les accents des cas frequents.
_ACRONYMES = {"devops", "rh", "ttm", "po", "kpi", "ci", "cd"}
_ACCENTS = {"a": "à", "agilite": "agilité", "echelle": "échelle",
            "equipe": "équipe", "qualite": "qualité", "strategie": "stratégie",
            "amelioration": "amélioration", "manageriale": "managériale",
            "modele": "modèle", "delivery": "delivery", "ingenierie": "ingénierie"}


def _nettoyer_label(texte):
    """Retire le contenu entre parentheses d'un libelle pilier/objectif (ex.
    « Ressources humaines (formations, coaching agile, talent, ...) » ->
    « Ressources humaines ») — mire `nettoyerLabel` de radar-svg.js. Les noms
    du referentiel Excel embarquent parfois un complement entre parentheses
    qui alourdit l'affichage sans aider a la lecture rapide ; on l'enleve
    PARTOUT ou joli_nom() est appele (barres, cartes, radar...), pas juste
    a un endroit — un oubli ponctuel a deja ete signale par le passe."""
    return re.sub(r"\s*\([^)]*\)", "", str(texte)).strip()


def joli_nom(nom):
    if not nom:
        return nom
    nom = _nettoyer_label(nom)
    mots = nom.replace("'", "' ").split()
    out = []
    premier = True
    for m in mots:
        bas = m.lower()
        if bas.endswith("'"):       # "l'" / "d'" : toujours minuscule, colle au suivant
            out.append(bas)
            continue
        coeur = _ACCENTS.get(bas, bas)
        if premier:                 # premier vrai mot : capitale
            coeur = coeur[:1].upper() + coeur[1:]
            premier = False
        elif bas in _ACRONYMES:     # DevOps & co restent en majuscules
            coeur = m.upper() if len(bas) <= 2 else coeur[:1].upper() + coeur[1:]
        out.append(coeur)
    return " ".join(out).replace("' ", "'")


def moyenne(valeurs):
    vals = [v for v in valeurs if v is not None]
    return sum(vals) / len(vals) if vals else None


def fleche(delta):
    if delta is None:
        return ("", D.MUTED)
    if abs(delta) < 0.05:
        return ("=", D.MUTED)
    return (f"▲ +{delta:.1f}", D.OK) if delta > 0 else (f"▼ {delta:.1f}", D.WARN)


def _trouver_layout(layouts, patterns, defaut_idx):
    """1er layout dont le nom contient un des `patterns` (insensible casse/accents
    grossiers) ; sinon le layout `defaut_idx` (borne). Permet de viser le bon
    layout meme si on fournit un autre template ou les indices different."""
    for lay in layouts:
        nom = (lay.name or "").lower()
        if any(p in nom for p in patterns):
            return lay
    return layouts[min(defaut_idx, len(layouts) - 1)]


def titre_slide(prs, layouts, texte):
    slide = prs.slides.add_slide(_trouver_layout(layouts, TITRE_PATTERNS, LAYOUT_TITRE_SEUL))
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            ph.text_frame.text = texte
            # Abaisse la taille du titre (le template la fige ~28pt) sans toucher
            # a la police/couleur du placeholder (Outfit navy) : on ne change que size.
            for para in ph.text_frame.paragraphs:
                for run in para.runs:
                    run.font.size = D.Pt(TITRE_SLIDE_PT)
            break
    return slide


def info_ligne(bloc):
    if bloc["type"] == "equipe":
        base = f"{bloc['effectif']} répondant(s)"
        return f"{bloc['departement']} · {base}" if bloc.get("departement") else base
    return f"{bloc['effectif']} répondant(s) · {bloc.get('nbEquipes', '')} équipe(s)"


# ----------------------------------------------------------------------------
# Slide 1 : Vue d'ensemble (jauge + barres par pilier + chips)
# ----------------------------------------------------------------------------
def _surtitre(slide, x, y, w, texte):
    """Petit label de section (capitales espacees) + filet fin dessous."""
    D.add_text(slide, x, y, w, 0.24,
               [(texte, {"size": D.TYPE["tiny"], "bold": True, "color": D.MUTED})])
    D.add_rect(slide, x, y + 0.26, w, 0.014, fill=CYAN)   # filet d'accent charte


def slide_vue_ensemble(prs, layouts, bloc):
    slide = titre_slide(prs, layouts, f"{bloc['nom']} — Vue d'ensemble")
    piliers = bloc.get("piliers", [])

    # Deltas par pilier (evolution), indexes par nom.
    deltas = {}
    comp = bloc.get("comparaison", {})
    if comp.get("disponible"):
        for p in comp.get("piliers", []):
            deltas[p["nom"]] = p.get("delta")

    glob = moyenne([p.get("moyenne") for p in piliers])
    glob_delta = None
    if comp.get("disponible"):
        cur = moyenne([p.get("courant") for p in comp.get("piliers", [])])
        prec = moyenne([p.get("precedent") for p in comp.get("piliers", [])])
        if cur is not None and prec is not None:
            glob_delta = cur - prec

    D.add_text(slide, MARGE_X, CONTENU_TOP - 0.10, 4.0, 0.3,
               [(info_ligne(bloc), {"size": D.TYPE["small"], "color": D.MUTED})])

    chips_top = CONTENU_BOTTOM - 0.55   # bandeau du bas (point fort / a renforcer)
    band_top = CONTENU_TOP + 0.42       # haut de la zone de contenu utile
    band_bot = chips_top - 0.30

    # --- Jauge moyenne globale (panneau gauche, centre verticalement) ---
    panel_w = 2.85
    D.add_rect(slide, MARGE_X, band_top, panel_w, band_bot - band_top,
               fill=FOND_PANNEAU, line=D.LINE, line_w=0.75, rounded=True, radius=RADIUS)
    cx = MARGE_X + panel_w / 2
    # Bloc {label + jauge + tendance} centre dans le panneau.
    has_delta = glob_delta is not None
    # 0.30 = label + gap jusqu'a la jauge (gy = top0 + 0.30) ; 2.0 = jauge ; 0.56 = ligne
    # de tendance quand comparaison. Colle au contenu reel pour un centrage exact
    # (evite le leger vide en bas sans comparaison, cf. revue design).
    bloc_h = 0.30 + 2.0 + (0.56 if has_delta else 0.0)
    top0 = band_top + (band_bot - band_top - bloc_h) / 2
    D.add_text(slide, MARGE_X, top0, panel_w, 0.24,
               [("MOYENNE GLOBALE", {"size": D.TYPE["tiny"], "bold": True, "color": D.MUTED,
                                     "align": PP_ALIGN.CENTER})])
    gs = 2.0
    gy = top0 + 0.30
    gx = cx - gs / 2
    frac = (glob / 3.0) if glob is not None else 0
    D.add_gauge(slide, gx, gy, gs, frac, ACCENT)
    D.add_text(slide, gx, gy, gs, gs,
               [(fmt(glob), {"size": D.TYPE["kpi"], "bold": True, "align": PP_ALIGN.CENTER}),
                ("sur 3", {"size": D.TYPE["small"], "color": D.MUTED, "align": PP_ALIGN.CENTER})],
               anchor=MSO_ANCHOR.MIDDLE)
    if has_delta:
        txt, col = fleche(glob_delta)
        D.add_text(slide, MARGE_X, gy + gs + 0.06, panel_w, 0.5,
                   [(txt, {"size": D.TYPE["h2"], "bold": True, "color": col,
                           "align": PP_ALIGN.CENTER}),
                    (f"vs {comp.get('precedenteDate', '')}",
                     {"size": D.TYPE["tiny"], "color": D.MUTED, "align": PP_ALIGN.CENTER})])

    # --- Barres par pilier (colonne droite) ---
    bx = MARGE_X + panel_w + 0.55
    bw = 10 - MARGE_X - bx
    _surtitre(slide, bx, CONTENU_TOP - 0.05, bw, "MATURITÉ PAR PILIER")
    n = max(1, len(piliers))
    label_w = 2.05
    val_w = 1.05
    track_x = bx + label_w
    track_w = bw - label_w - val_w
    zone_top = CONTENU_TOP + 0.42
    zone_h = band_bot - zone_top
    row_h = min(0.66, zone_h / n)
    y = zone_top + (zone_h - row_h * n) / 2   # centre le paquet de barres
    bar_h = 0.24
    for i, p in enumerate(piliers):
        moy = p.get("moyenne")
        col_pil = D.couleur_pilier(i)
        D.add_dot(slide, bx, y + (row_h - 0.12) / 2, 0.12, col_pil)
        D.add_text(slide, bx + 0.22, y, label_w - 0.22, row_h,
                   [(joli_nom(p["nom"]), {"size": D.TYPE["small"], "bold": True,
                                          "line_spacing": 0.95})],
                   anchor=MSO_ANCHOR.MIDDLE)
        by = y + (row_h - bar_h) / 2
        D.add_hbar(slide, track_x, by, track_w, bar_h,
                   (moy / 3.0) if moy is not None else 0, col_pil)
        txt, col = fleche(deltas.get(p["nom"]))
        lignes = [(fmt(moy), {"size": D.TYPE["h3"], "bold": True})]
        if txt:
            lignes.append((txt, {"size": D.TYPE["tiny"], "bold": True, "color": col}))
        D.add_text(slide, track_x + track_w + 0.12, y, val_w - 0.12, row_h,
                   lignes, anchor=MSO_ANCHOR.MIDDLE)
        y += row_h
    # Reglette d'echelle 0-1-2-3 sous les barres.
    sy = zone_top + (zone_h + row_h * n) / 2 - 0.02
    for g in range(4):
        gx2 = track_x + track_w * (g / 3.0)
        D.add_rect(slide, gx2 - 0.005, zone_top + (zone_h - row_h * n) / 2,
                   0.01, row_h * n, fill=D.LINE)
        D.add_text(slide, gx2 - 0.15, sy, 0.30, 0.18,
                   [(str(g), {"size": D.TYPE["tiny"], "color": D.MUTED,
                              "align": PP_ALIGN.CENTER})])

    # --- Bandeau bas : point fort / a renforcer ---
    valides = [(i, p["nom"], p["moyenne"]) for i, p in enumerate(piliers) if p.get("moyenne") is not None]
    if valides:
        fort = max(valides, key=lambda x: x[2])
        faible = min(valides, key=lambda x: x[2])
        band_w = BORD_DROIT - MARGE_X            # s'arrete avant le badge n° de slide
        D.add_rect(slide, MARGE_X, chips_top, band_w, 0.52,
                   fill=FOND_PANNEAU, line=D.LINE, line_w=0.75, rounded=True, radius=RADIUS)
        mid = MARGE_X + band_w / 2               # divise le bandeau en deux moities egales
        _chip(slide, MARGE_X + 0.30, chips_top, mid - MARGE_X - 0.30, "Point fort", "▲", fort)
        D.add_rect(slide, mid, chips_top + 0.10, 0.012, 0.32, fill=D.LINE)
        _chip(slide, mid + 0.30, chips_top, BORD_DROIT - mid - 0.30, "À renforcer", "▼", faible)


def _chip(slide, x, y, w, prefixe, glyphe, pilier):
    # Pastille = couleur d'IDENTITE du pilier nomme (mire sa barre au-dessus) ; le
    # SENS (force/faiblesse) est porte par le glyphe ▲/▼, plus par la couleur.
    idx, nom, moy = pilier
    D.add_dot(slide, x, y + 0.20, 0.13, D.couleur_pilier(idx))
    box = slide.shapes.add_textbox(Inches(x + 0.24), Inches(y), Inches(w - 0.24), Inches(0.52))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.auto_size = None
    for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(tf, m, 0)
    p = tf.paragraphs[0]
    rg = p.add_run(); rg.text = f"{glyphe} "
    rg.font.size = D.Pt(D.TYPE["small"]); rg.font.bold = True; rg.font.color.rgb = D.rgb(D.INK)
    r1 = p.add_run(); r1.text = f"{prefixe} : "
    r1.font.size = D.Pt(D.TYPE["small"]); r1.font.color.rgb = D.rgb(D.MUTED)
    r2 = p.add_run(); r2.text = f"{joli_nom(nom)} — {fmt(moy)} / 3"
    r2.font.size = D.Pt(D.TYPE["small"]); r2.font.bold = True; r2.font.color.rgb = D.rgb(D.INK)
    if D.POLICE:                        # meme police de marque que le reste du deck
        for r in (rg, r1, r2):
            r.font.name = D.POLICE


# ----------------------------------------------------------------------------
# Slide 2 : Radar par objectif + commentaire + evolution
# ----------------------------------------------------------------------------
# Radar VECTORIEL (formes python-pptx natives, pas un PNG rasterise par
# Puppeteer) : reste net a toute resolution/impression et reste editable dans
# PowerPoint. Mire le meme principe visuel que le radar web (voir
# app/src/radar-svg.js, qui garde son propre role serveur/US6.2-US6.5 — on ne
# le touche pas) : grille de niveaux + rayons, polygone "precedent" pointille
# (comparaison), polygone "courant" en aire semi-transparente, une puce
# couleur-pilier par sommet, libelles d'axe colores par pilier, legende a
# droite. Echelle fixe 0..3 (les 4 niveaux de reponse du referentiel).
RADAR_MAX = 3
# Cote max du cercle : le panneau de DROITE (commentaire + evolution par pilier)
# se partage la largeur de slide et a besoin de place pour ses propres noms de pilier.
RADAR_COTE_MAX = 4.2
# Taille LISIBLE des libelles d'axe, propre au radar (pas soumise au ×0.9 global) :
# au 0.9 ils tombaient a ~7pt, illisibles (retour utilisateur 2026-07-22). Le cercle
# ayant de la place depuis que la legende est passee en bandeau horizontal au-dessus,
# les libelles peuvent respirer plus gros. _taille_libelle_axe reduit dans cette borne
# si un mot compose long l'exige, puis _forcer_cesure prend le relai.
RADAR_LABEL_MAX = 9.0
RADAR_LABEL_MIN = 8.0
# Hauteur de ligne REELLE (pas juste la taille de police) : mesuree au rendu,
# comme LH_QUESTION (0.195 pour small 10.5pt) — proportionnelle a la taille
# pour tiny (9pt). Sous-estimer ceci fait deborder le texte de sa boite sans
# que le controle geometrique ne le voie (la FORME reste dans le cadre).
RADAR_LH = 0.195 * 9 / 10.5   # ~0.167, hauteur de ligne tiny (9pt)
# Bandeau de section (_surtitre, meme grammaire que "MATURITÉ PAR PILIER" sur la
# vue d'ensemble), reserve AU-DESSUS du cercle. (La reglette d'echelle horizontale
# 0-3 a ete RETIREE le 2026-07-22 : trompeuse sur un radar radial, cf. revue design.)
RADAR_HEADER_H = 0.42


def _cote_radar(w, h):
    """Cote (in) du carre du radar pour une boite (w, h) donnee — factorise
    pour que slide_radar et _dessiner_radar (formes) calculent EXACTEMENT la
    meme valeur. La legende couleurs/series ayant ete deportee EN BANDEAU
    HORIZONTAL au-dessus du cercle (revue design 2026-07-22), la boite n'a plus
    a reserver de colonne de legende a droite : le cercle occupe toute sa boite."""
    return min(h, RADAR_COTE_MAX, w)


def _point_radar(cx, cy, rayon, i, n, valeur):
    ang = -math.pi / 2 + i * (2 * math.pi / n)
    r = rayon * max(0.0, min(RADAR_MAX, valeur if valeur is not None else 0)) / RADAR_MAX
    return (cx + r * math.cos(ang), cy + r * math.sin(ang))


def _lignes_radar(texte, largeur_in, taille=None):
    taille = taille if taille is not None else D.TYPE["tiny"]
    return max(1, D.estimer_lignes(texte, max(0.3, largeur_in), taille))


def _taille_libelle_axe(nom_axe, box_w, taille_max=None, taille_min=7.0,
                        cpi_ref=11.0, taille_ref=10.5):
    """Reduit la taille de police jusqu'a ce que le mot le plus long de
    `nom_axe` (separe sur espace ET tiret — un tiret existant est deja un
    point de coupure sur) tienne dans `box_w`. Necessaire car D.estimer_lignes
    (repli mot-a-mot par nombre de caracteres) sous-estime le nombre de lignes
    reel quand un seul mot compose francais ("Fonctionnement",
    "Synchronisation") depasse a lui seul la largeur de la boite : PowerPoint
    le coupe alors au milieu, SANS trait d'union, sur un radar dense (10-12
    axes, boites de libelle etroites). Repli sur `taille_min` si meme a cette
    taille le mot ne tient pas (mot exceptionnellement long — voir
    `_forcer_cesure`, appele par l'appelant dans ce cas)."""
    taille_max = taille_max if taille_max is not None else D.TYPE["tiny"]
    mots = [m for m in re.split(r"[ \-]", nom_axe) if m]
    plus_long = max((len(m) for m in mots), default=0)
    taille = taille_max
    while taille > taille_min:
        cpi = cpi_ref * (taille_ref / taille)
        if plus_long <= box_w * cpi:
            break
        taille = round(taille - 0.5, 2)
    return max(taille, taille_min)


def _forcer_cesure(nom_axe, box_w, taille, cpi_ref=11.0, taille_ref=10.5,
                   min_prefixe=3):
    """Insere un vrai trait d'union dans tout mot de `nom_axe` encore trop
    long pour `box_w` a `taille` (meme apres le plancher de
    `_taille_libelle_axe`) — au point ou PowerPoint le couperait de toute
    facon, mais avec un tiret plutot qu'une coupure brute. Un mot deja
    cesure ("inter-équipes") n'est jamais retouche (deja un point de coupure
    propre). Sans effet si tous les mots tiennent deja."""
    cpi = cpi_ref * (taille_ref / taille)
    max_chars = max(min_prefixe + 1, int(box_w * cpi))

    def _cesurer_mot(mot):
        if len(mot) <= max_chars or "-" in mot:
            return mot
        coupe = max(min_prefixe, max_chars - 1)
        coupe = min(coupe, len(mot) - 1)
        return mot[:coupe] + "-" + mot[coupe:]

    return " ".join(_cesurer_mot(m) for m in nom_axe.split(" "))


def _dessiner_radar(slide, x, y, w, h, axes):
    """Dessine le radar (grille + polygones + libelles d'axe) dans la boite
    (x, y, w, h). `axes` = [{nom, moyenne, precedent, pilierIndex}]. La legende
    couleurs/series est dessinee SEPAREMENT (cf. _legende_radar_verticale, colonne
    a droite du cercle sur la slide radar dediee) : ici le cercle est centre dans
    sa boite et les libelles d'axe s'etendent dans les marges. Les libelles sont
    dimensionnes a leur contenu reel pour ne jamais se chevaucher."""
    n = len(axes)
    if n < 3:
        return  # radar illisible sous 3 axes : rien plutot qu'une forme deformee
    cote = _cote_radar(w, h)
    x0, y0 = x + (w - cote) / 2, y + (h - cote) / 2   # cercle centre dans sa boite
    cx, cy = x0 + cote / 2, y0 + cote / 2
    rayon = cote * 0.37   # grille agrandie (slide radar dediee) : les libelles ont la largeur pour respirer
    rlim = x + w   # bord droit de la boite = limite des libelles cote droit (plus de legende)

    has_prev = any(a.get("precedent") is not None for a in axes)

    # Grille (niveaux 1..MAX) + rayons, ton neutre discret (theme).
    for niveau in range(1, RADAR_MAX + 1):
        pts = [_point_radar(cx, cy, rayon, i, n, niveau) for i in range(n)]
        D.add_polygon(slide, pts, line=D.LINE, line_w=0.75)
    for i in range(n):
        px_, py_ = _point_radar(cx, cy, rayon, i, n, RADAR_MAX)
        D.add_line(slide, cx, cy, px_, py_, D.LINE, width=0.75)

    # Polygone "session precedente" (pointille, sans remplissage). Un axe sans
    # comparaison reprend sa valeur courante (evite un effondrement a 0 qui
    # laisserait croire a une regression totale la ou il n'y a simplement pas
    # de comparaison possible) — mire le meme choix que radar-svg.js.
    if has_prev:
        pts_prec = [_point_radar(cx, cy, rayon, i, n,
                                 a.get("precedent") if a.get("precedent") is not None
                                 else a.get("moyenne")) for i, a in enumerate(axes)]
        D.add_polygon(slide, pts_prec, line=D.MUTED, line_w=1.5, dash=D.DASH.DASH)

    # Polygone "courant" : aire semi-transparente, couleur neutre du radar
    # (palette pilier n°0 — comme radar-svg.js, ce n'est pas une couleur de
    # marque) ; la couleur de CHAQUE pilier vit sur sa puce et son libelle.
    couleur_aire = D.couleur_pilier(0)
    pts_cour = [_point_radar(cx, cy, rayon, i, n, a.get("moyenne")) for i, a in enumerate(axes)]
    D.add_polygon(slide, pts_cour, fill=couleur_aire, alpha=27, line=couleur_aire, line_w=2)
    for i, a in enumerate(axes):
        px_, py_ = pts_cour[i]
        d = max(0.06, cote * 0.014)
        D.add_dot(slide, px_ - d / 2, py_ - d / 2, d, D.couleur_pilier(a.get("pilierIndex", 0)))

    # Libelles d'axe : petite zone de texte centree sur le sommet exterieur,
    # alignee selon le cote du radar (gauche/centre/droite selon le signe du
    # cosinus) — mire le text-anchor start/middle/end du radar web. Largeur
    # et hauteur de boite calculees (pas devinees) pour ne jamais chevaucher
    # ni la colonne de legende (a droite) ni le bord de la zone (a gauche).
    # Cap a 3 lignes (ellipse au-dela, D.tronquer_a_lignes) : sur un radar dense
    # (10-12 axes), un libelle demesurement long empieterait sinon sur les
    # libelles voisins (peu d'espace vertical entre sommets adjacents). Le
    # budget de hauteur ajoute une demi-ligne de marge : l'estimateur de repli
    # mot-a-mot (D.estimer_lignes) sous-estime le nombre de lignes reel quand
    # un seul mot (compose, frequent en francais — "Fonctionnement",
    # "Synchronisation") depasse a lui seul la largeur de la boite.
    # Largeur de libelle genereuse (2.1) : le radar occupe desormais une slide dediee
    # pleine largeur -> de larges marges laterales pour etaler les libelles d'axe sans
    # ellipse ni collision (meme les longs comme "Orientation client et pilotage...").
    LARGEUR_MAX_LABEL, MAX_LIGNES_LABEL = 2.1, 3
    for i, a in enumerate(axes):
        ang = -math.pi / 2 + i * (2 * math.pi / n)
        cosang = math.cos(ang)
        # Libelles pousses un peu plus loin du cercle (0.075 vs 0.05) : degage les
        # sommets voisins en haut/bas d'un radar dense (12 axes) ou les boites se touchent.
        lx = cx + (rayon + cote * 0.075) * cosang
        ly = cy + (rayon + cote * 0.075) * math.sin(ang)
        nom_axe = joli_nom(a.get("nom", ""))
        if cosang > 0.2:
            box_w = max(0.65, min(LARGEUR_MAX_LABEL, rlim - lx - 0.08))
            box_x, align = lx, PP_ALIGN.LEFT
        elif cosang < -0.2:
            box_w = max(0.65, min(LARGEUR_MAX_LABEL, lx - x - 0.08))
            box_x, align = lx - box_w, PP_ALIGN.RIGHT
        else:
            box_w = min(LARGEUR_MAX_LABEL, cote)
            box_x, align = lx - box_w / 2, PP_ALIGN.CENTER
        box_x = max(0.02, min(box_x, 10 - box_w - 0.02))
        # Reduit la taille AVANT le repli mot-a-mot si un mot seul ne tiendrait
        # pas dans box_w a la taille normale (cf. note _taille_libelle_axe) —
        # evite que PowerPoint coupe ce mot au milieu sans trait d'union.
        # Les libelles d'axe ont une taille LISIBLE propre (RADAR_LABEL_*), non
        # soumise au ×0.9 global : au 0.9 ils devenaient illisibles (~7pt), et le
        # cercle a de la place depuis que la legende est passee en bandeau au-dessus.
        taille_axe = _taille_libelle_axe(nom_axe, box_w, taille_max=RADAR_LABEL_MAX,
                                         taille_min=RADAR_LABEL_MIN)
        nom_axe = _forcer_cesure(nom_axe, box_w, taille_axe)
        lh_axe = RADAR_LH * (taille_axe / D.TYPE["tiny"])
        nom_axe = D.tronquer_a_lignes(nom_axe, box_w, taille_axe, MAX_LIGNES_LABEL)
        box_h = max(0.20, (min(MAX_LIGNES_LABEL, _lignes_radar(nom_axe, box_w, taille_axe)) + 0.5)
                    * lh_axe + 0.06)
        # Pastille pilier sur l'axe, juste avant le libellé : la couleur passe du
        # TEXTE a la pastille — libellés en D.INK (lisibilité + contraste WCAG ; le
        # gold #b8860b en texte échouait a 3.25:1), non gras (moins lourd, moins de
        # collisions). Mire le radar web (pastille + texte foncé, comme la légende).
        dd = max(0.06, cote * 0.013)
        drx = cx + (rayon + cote * 0.012) * cosang   # juste hors du cercle : degage le libelle (a rayon+0.05*cote)
        dry = cy + (rayon + cote * 0.012) * math.sin(ang)
        D.add_dot(slide, drx - dd / 2, dry - dd / 2, dd, D.couleur_pilier(a.get("pilierIndex", 0)))
        D.add_text(slide, box_x, ly - box_h / 2, box_w, box_h,
                   [(nom_axe, {"size": taille_axe, "bold": False,
                      "color": D.INK,
                      "align": align, "line_spacing": 0.95})],
                   anchor=MSO_ANCHOR.MIDDLE, align=align)


def _legende_radar_verticale(slide, x, y, w, h, piliers, has_prev, couleur_aire):
    """Légende du radar en colonne VERTICALE, compacte, calée à droite du cercle et
    CENTRÉE verticalement dans la hauteur `h` (demande utilisateur 2026-07-22 : légende
    à droite en plus petit, radar recentré/agrandi). Puce couleur + nom de pilier par
    ligne, puis les deux clés de série (Session courante/précédente). Chaque ligne est
    dimensionnée à son propre nombre de lignes réel (un nom court n'hérite pas de la
    hauteur d'un nom long)."""
    size = D.TYPE["tiny"]
    dot = 0.12
    rows = [(i, joli_nom(nom),
             max(0.24, D.estimer_lignes(joli_nom(nom), w - 0.24, size) * RADAR_LH + 0.08))
            for i, nom in enumerate(piliers)]
    comp_h = (0.14 + 2 * 0.24) if has_prev else 0.0
    total = sum(rh for _, _, rh in rows) + comp_h
    if total > h:   # garde-fou : jamais deborder (beaucoup de piliers)
        f = h / total
        rows = [(i, t, rh * f) for i, t, rh in rows]
        comp_h *= f
        total = h
    cy = y + max(0.0, (h - total) / 2)   # centrage vertical
    for i, txt, rh in rows:
        D.add_dot(slide, x, cy + (min(rh, 0.28) - dot) / 2, dot, D.couleur_pilier(i))
        D.add_text(slide, x + 0.24, cy, w - 0.24, rh,
                   [(txt, {"size": size, "line_spacing": 1.0})], anchor=MSO_ANCHOR.MIDDLE)
        cy += rh
    if has_prev:
        cy += 0.12
        D.add_line(slide, x, cy + 0.10, x + 0.28, cy + 0.10, couleur_aire, width=2.5)
        D.add_text(slide, x + 0.34, cy, w - 0.34, 0.22,
                   [("Session courante", {"size": size, "color": D.MUTED})],
                   anchor=MSO_ANCHOR.MIDDLE)
        cy += 0.24
        D.add_line(slide, x, cy + 0.10, x + 0.28, cy + 0.10, D.MUTED, width=2.5,
                   dash=D.DASH.DASH)
        D.add_text(slide, x + 0.34, cy, w - 0.34, 0.22,
                   [("Session précédente", {"size": size, "color": D.MUTED})],
                   anchor=MSO_ANCHOR.MIDDLE)


def _hauteur_commentaire(texte, largeur_in):
    """Estime la hauteur (in) d'un commentaire wrappe en small, pour eviter le
    grand vide d'un encart a hauteur fixe. On simule le wrap mot-a-mot avec une
    largeur prudente (~12.5 car/inch, calibree pour ce corps de texte) et on
    majore un peu, pour ne JAMAIS sous-estimer (un debordement est pire qu'un
    petit blanc)."""
    if not texte:
        return 0.60
    lignes = D.estimer_lignes(texte, largeur_in, D.TYPE["small"], cpi_ref=12.5,
                              taille_ref=D.TYPE["small"])
    # entete (~0.30) + lignes (small 10.5pt + interligne ~0.205in) + marges
    return 0.42 + lignes * 0.205


def slide_radar(prs, layouts, bloc):
    """Slide DÉDIÉE au radar (scindée du commentaire/progression le 2026-07-22,
    demande utilisateur) : le radar occupe toute la largeur de contenu -> cercle
    centré et libellés d'axe qui s'étalent dans de larges marges (plus d'ellipse,
    pas de collision). Le commentaire et la progression vivent sur slide_progression,
    juste après. Vectoriel (pptx_deck.add_polygon/add_line) à partir des mêmes
    données (objectifs/piliers) que le PNG serveur — net à toute résolution, éditable."""
    slide = titre_slide(prs, layouts, f"{bloc['nom']} — Radar de maturité")
    axes = bloc.get("objectifs") or []
    piliers_legende = [p.get("nom", "") for p in (bloc.get("piliers") or [])]
    W = 10 - 2 * MARGE_X   # pleine largeur de contenu
    if len(axes) < 3:
        # Radar illisible sous 3 axes : message plutot qu'une forme deformee.
        _surtitre(slide, MARGE_X, CONTENU_TOP, W, "MATURITÉ PAR OBJECTIF")
        D.add_text(slide, MARGE_X, CONTENU_TOP + 0.9, W, 0.4,
                   [("Radar indisponible (moins de 3 objectifs mesurés).",
                     {"size": D.TYPE["small"], "italic": True, "color": D.MUTED})])
        return
    # Légende couleurs/séries VERTICALE, compacte, à DROITE (demande utilisateur
    # 2026-07-22) : le radar récupère toute la HAUTEUR de contenu à gauche -> cercle
    # bien plus grand et recentré dans sa zone. Pas de sur-titre (redondant avec le
    # titre de slide "Radar de maturité").
    H = CONTENU_BOTTOM - CONTENU_TOP
    couleur_aire = D.couleur_pilier(0)
    has_prev = any(a.get("precedent") is not None for a in axes)
    LEGENDE_W, GAP_LEG = 1.95, 0.30
    radar_w = W - LEGENDE_W - GAP_LEG
    _dessiner_radar(slide, MARGE_X, CONTENU_TOP, radar_w, H, axes)
    _legende_radar_verticale(slide, MARGE_X + radar_w + GAP_LEG, CONTENU_TOP, LEGENDE_W, H,
                             piliers_legende, has_prev, couleur_aire)


def slide_progression(prs, layouts, bloc):
    """Slide « Progression & commentaire » (scindée du radar le 2026-07-22) :
    le commentaire de restitution (callout pleine largeur) puis l'évolution par
    pilier (précédent -> courant, barre + delta). Disposée sur toute la largeur,
    plus lisible que l'ancienne colonne étroite à droite du radar."""
    slide = titre_slide(prs, layouts, f"{bloc['nom']} — Progression & commentaire")
    x = MARGE_X
    W = 10 - 2 * MARGE_X
    commentaire = (bloc.get("commentaire") or "").strip()
    comp = bloc.get("comparaison", {})
    piliers_ev = comp.get("piliers", []) if comp.get("disponible") else []

    # ---- Commentaire de restitution : callout pleine largeur, en haut ----
    txt_in = W - 0.52
    h_comm = max(1.05, _hauteur_commentaire(commentaire, txt_in) + 0.22)
    top = CONTENU_TOP
    D.add_rect(slide, x, top, W, h_comm, fill=FOND_PANNEAU, line=D.LINE,
               line_w=0.75, rounded=True, radius=RADIUS)
    D.add_rect(slide, x, top, 0.07, h_comm, fill=CYAN, rounded=True, radius=0.5)
    lignes = [("COMMENTAIRE DE RESTITUTION",
               {"size": D.TYPE["tiny"], "bold": True, "color": D.MUTED, "space_after": 5})]
    if commentaire:
        for ligne in commentaire.split("\n"):
            lignes.append((ligne, {"size": D.TYPE["small"], "space_after": 4,
                                   "line_spacing": 1.06}))
    else:
        lignes.append(("(à compléter)", {"size": D.TYPE["small"], "italic": True,
                                         "color": D.MUTED}))
    D.add_text(slide, x + 0.28, top + 0.16, txt_in, h_comm - 0.30, lignes,
               anchor=MSO_ANCHOR.TOP)

    # ---- Évolution par pilier : nom · barre (courant, avec repère précédent) · av→ap · delta ----
    if not piliers_ev:
        return
    ey = top + h_comm + 0.34
    _surtitre(slide, x, ey, W, f"ÉVOLUTION VS {comp.get('precedenteDate', '')}".upper())
    rows_top = ey + 0.46
    bottom = CONTENU_BOTTOM - 0.14
    n = len(piliers_ev)
    # Une ligne par pilier, hauteur egale repartie dans l'espace restant (borne
    # pour rester lisible/aere sans etirer demesurement si peu de piliers).
    row_h = max(0.44, min(0.78, (bottom - rows_top) / max(1, n)))
    # Colonnes : nom (gauche) | barre 0-3 | "av -> ap" | delta (droite), calees a droite
    # en s'arretant avant le badge n° de slide.
    right_lim = x + W - 0.20
    delta_w, avap_w = 0.85, 1.05
    delta_x = right_lim - delta_w
    avap_x = delta_x - avap_w
    bar_w = 2.4
    bar_x = avap_x - 0.30 - bar_w
    name_w = max(1.2, bar_x - 0.30 - (x + 0.24))
    y = rows_top
    for i, p in enumerate(piliers_ev):
        cy = y + row_h / 2
        D.add_dot(slide, x, cy - 0.06, 0.12, D.couleur_pilier(i))
        D.add_text(slide, x + 0.24, y, name_w, row_h,
                   [(joli_nom(p["nom"]), {"size": D.TYPE["small"], "bold": True,
                                          "line_spacing": 0.95})],
                   anchor=MSO_ANCHOR.MIDDLE)
        cour = p.get("courant")
        prec = p.get("precedent")
        by = cy - 0.065
        D.add_hbar(slide, bar_x, by, bar_w, 0.13,
                   (cour / 3.0) if cour is not None else 0, D.couleur_pilier(i))
        if prec is not None:   # repere de la valeur precedente sur la piste
            mx = bar_x + bar_w * max(0.0, min(1.0, prec / 3.0))
            D.add_rect(slide, mx - 0.008, by - 0.05, 0.016, 0.23, fill=D.INK)
        D.add_text(slide, avap_x, y, avap_w, row_h,
                   [(f"{fmt(prec)} → {fmt(cour)}",
                     {"size": D.TYPE["small"], "color": D.MUTED, "align": PP_ALIGN.RIGHT})],
                   anchor=MSO_ANCHOR.MIDDLE)
        txt, col = fleche(p.get("delta"))
        if txt:
            D.add_text(slide, delta_x, y, delta_w, row_h,
                       [(txt, {"size": D.TYPE["small"], "bold": True, "color": col,
                               "align": PP_ALIGN.RIGHT})],
                       anchor=MSO_ANCHOR.MIDDLE)
        y += row_h


# ----------------------------------------------------------------------------
# Slide 3 : Points d'attention (cartes)
# ----------------------------------------------------------------------------
def _entete_colonne(slide, x, w, glyphe, titre, sous):
    """En-tete de colonne epure (meme grammaire que _surtitre / "MATURITÉ PAR
    PILIER") : glyphe de SENS (▲ positif / ▼ à travailler) en navy SANS fond +
    titre + sous-titre, cloture par un filet CYAN (accent de charte). La couleur
    n'encode toujours pas le sens (force/faiblesse) — reservee a l'identite des
    piliers ; le sens passe par le glyphe, le filet est un pur accent de marque."""
    D.add_text(slide, x, CONTENU_TOP - 0.10, 0.28, 0.34,
               [(glyphe, {"size": D.TYPE["h3"], "bold": True, "color": D.INK})],
               anchor=MSO_ANCHOR.MIDDLE)
    D.add_text(slide, x + 0.30, CONTENU_TOP - 0.10, w - 0.30, 0.34,
               [(titre, {"size": D.TYPE["h3"], "bold": True})])
    D.add_text(slide, x + 0.30, CONTENU_TOP + 0.26, w - 0.30, 0.22,
               [(sous, {"size": D.TYPE["tiny"], "color": D.MUTED})])
    D.add_rect(slide, x, CONTENU_TOP + 0.52, w, 0.014, fill=CYAN)   # filet d'accent charte


# Taille plancher des cartes de points forts/attention (voir D.ajuster_police).
# En-dessous, le filet de securite de _cartes_colonne prend le relai (troncature
# avec ellipse) plutot que de laisser une forme deborder de la slide.
TAILLE_MIN_CARTE = 6.0
GAP_MIN, GAP_MAX = 0.14, 0.28
# Nombre de cartes par colonne. 2 (2026-07-22, arbitrage utilisateur) : a 3 cartes, 3
# questions longues + le chrome fixe (barre/contexte/marges) ne tenaient qu'au plancher
# 6pt (illisible) ; a 2 cartes chaque question respire et prend une taille lisible. On
# montre donc le top 2 par colonne.
N_CARTES_MAX = 2


# Hauteur d'une carte HORS question : gap(0.04) + contexte(0.17) + gap(0.10) +
# barre + label "moy."(0.41). Source unique partagee entre _bloc_carte_h (hauteur) et
# le calcul inverse de ql_max dans _cartes_colonne — evite qu'un des deux derive de
# l'autre. Le budget "barre" inclut le label "moy. X.X" pose sous le repere de moyenne
# des widgets de dispersion (les cartes a barre simple gardent un peu de marge en bas).
_CARTE_H_FIXE = 0.04 + 0.17 + 0.10 + 0.41   # = 0.72 : dimensionne TOUTES les cartes (budget commun)
# Contenu REEL d'une carte "score" (barre + "sur 3", SANS ligne "moy.") : sert
# uniquement a RE-CENTRER ces cartes (sinon, centrees sur le budget commun 0.72, elles
# laissent un vide en bas la ou les cartes "dispersion" mettent leur "moy."). Ne change
# ni le dimensionnement ni ql_max, donc pas la troncature.
_CARTE_H_FIXE_SCORES = 0.04 + 0.17 + 0.10 + 0.28   # = 0.59


# Hauteur du contenu d'une carte = question (ql lignes, a `taille` pt) + contexte
# + barre. Seule la hauteur de la question depend de `taille` (contexte et barre
# gardent une taille fixe) — les rendus placent leurs elements aux memes offsets,
# d'ou cette source unique.
def _bloc_carte_h(ql, taille=D.TYPE["small"], fixe=_CARTE_H_FIXE):
    lh = LH_QUESTION * (taille / D.TYPE["small"])
    return ql * lh + fixe


def _texte_et_lignes(texte, tw, taille, ql_max):
    """Nombre de lignes reellement utilisees (borne a `ql_max`, cf. filet de
    securite de _cartes_colonne) et texte eventuellement tronque (ellipse) si
    le texte naturel depasse cette borne."""
    ql_naturel = D.estimer_lignes(texte, tw, taille)
    ql = min(ql_naturel, ql_max)
    if ql < ql_naturel:
        texte = D.tronquer_a_lignes(texte, tw, taille, ql)
    return texte, ql


def _contexte_joli(ctx, max_chars=56):
    """Embellit le contexte "Pilier · Objectif" (le pilier vient de l'Excel en
    casse maladroite) et le borne a une ligne (tronque l'objectif au besoin)."""
    if not ctx:
        return ""
    parts = ctx.split(" · ")
    parts[0] = joli_nom(parts[0])
    txt = " · ".join(parts)
    if len(txt) > max_chars:
        txt = txt[:max_chars - 1].rstrip(" ,;(") + "…"
    return txt


def _valeur_cote_barre(slide, x, ry, w, lignes):
    """Pose le couple de valeur (2 lignes) a droite d'une barre (h=0.13), centre
    verticalement SUR la barre. Evite que le grand chiffre flotte au-dessus."""
    box_h = 0.42
    D.add_text(slide, x, ry + 0.065 - box_h / 2, w, box_h, lignes,
               anchor=MSO_ANCHOR.MIDDLE)


def _label_moyenne(slide, tx, rw, ry, moy):
    """Nomme le repere de moyenne pose par add_range_bar : sans lui, on lit un point
    non explique a cote d'un "ecart-type" (et quand la dispersion est nulle, l'amplitude
    est invisible, il ne reste QUE ce point). Pose "moy. X.X" sous le repere (meme
    abscisse : tx + rw * moy/3), borne a la largeur de la barre."""
    if moy is None:
        return
    box_w = 0.85
    mxx = tx + rw * max(0.0, min(1.0, moy / 3.0))
    box_x = max(tx, min(mxx - box_w / 2, tx + rw - box_w))
    D.add_text(slide, box_x, ry + 0.19, box_w, 0.16,
               [(f"moy. {fmt(moy)}", {"size": D.TYPE["tiny"], "color": D.MUTED,
                                      "align": PP_ALIGN.CENTER})],
               align=PP_ALIGN.CENTER)


def _widget_amplitude(slide, tx, ry, rw, mn, mx, moy):
    """Widget de dispersion mutualise (cartes accords/desaccords) : barre d'amplitude
    min–max slate + repere de moyenne navy + label 'moy. X.X'. Cas dispersion NULLE
    (min==max) : une pastille pleine 'consensus' a la moyenne + mention explicite,
    plutot qu'un repere isole sur une piste vide (qui se lit comme une barre cassee)."""
    consensus = moy is not None and mn is not None and mx is not None and (mx - mn) < 1e-6
    D.add_range_bar(slide, tx, ry, rw, 0.13, mn, mx, 3.0, D.MUTED,
                    marker=None if consensus else moy)
    if consensus:
        px = tx + rw * max(0.0, min(1.0, moy / 3.0))
        pill_w = 0.44
        D.add_rect(slide, max(tx, min(px - pill_w / 2, tx + rw - pill_w)), ry, pill_w, 0.13,
                   fill=D.INK, rounded=True, radius=0.5)
        D.add_text(slide, max(tx, min(px - 0.55, tx + rw - 1.1)), ry + 0.19, 1.1, 0.16,
                   [(f"moy. {fmt(moy)} · consensus",
                     {"size": D.TYPE["tiny"], "color": D.MUTED, "align": PP_ALIGN.CENTER})],
                   align=PP_ALIGN.CENTER)
    else:
        _label_moyenne(slide, tx, rw, ry, moy)


def _taille_colonne(items, w):
    """Taille de police qu'une colonne de cartes prendrait seule : la plus grande
    entre `small` et TAILLE_MIN_CARTE telle que la pile des cartes tienne dans la
    bande. Isolee de _cartes_colonne pour pouvoir calculer une taille COMMUNE a
    plusieurs colonnes (cf. _taille_cartes_bloc)."""
    items = items[:N_CARTES_MAX]
    n = max(1, len(items))
    band = CONTENU_BOTTOM - (CONTENU_TOP + 0.66)
    tw_estim = w - 0.42
    textes = [it.get("texte", "") for it in items]

    def budget_ok(taille, _lignes_max):
        gap = GAP_MIN if n > 1 else 0.0
        hs = [_bloc_carte_h(D.estimer_lignes(t, tw_estim, taille), taille) + 0.26 for t in textes]
        return sum(hs) + max(0, n - 1) * gap <= band

    taille, _ = D.ajuster_police(textes, tw_estim, D.TYPE["small"], TAILLE_MIN_CARTE, budget_ok)
    return taille


def _taille_cartes_bloc(bloc, colw):
    """Taille de carte COMMUNE aux slides Points forts (4) et Points d'attention (5)
    d'un bloc (demande utilisateur 2026-07-22 : les deux slides doivent afficher la
    MEME taille de texte). Prend le min des tailles que chaque colonne prendrait
    seule — le min garantit que les 4 colonnes tiennent a cette taille unique."""
    colonnes = [bloc.get("hauts", []), bloc.get("accords", []),
                bloc.get("dispersion", []), bloc.get("faibles", [])]
    tailles = [_taille_colonne(c, colw) for c in colonnes if c]
    return min(tailles) if tailles else D.TYPE["small"]


def _cartes_colonne(slide, x, w, items, accent, rendu, taille_forcee=None):
    items = items[:N_CARTES_MAX]   # top N par colonne (cf. N_CARTES_MAX)
    n = max(1, len(items))
    top = CONTENU_TOP + 0.66
    band = CONTENU_BOTTOM - top
    PAD_CARTE = 0.26
    # Largeur utile du texte dans la carte = w - pad(0.24) - marge droite(0.18),
    # cf. les rendus (tx = x + pad, tw = w - pad - 0.18).
    tw_estim = w - 0.42
    textes = [it.get("texte", "") for it in items]

    # Hauteur de CHAQUE carte a `taille` (pas un maximum commun applique a
    # toutes) : une question courte a droite d'une question tres longue ne doit
    # pas heriter de la hauteur de cette derniere.
    def hauteurs(taille):
        return [_bloc_carte_h(D.estimer_lignes(t, tw_estim, taille), taille) + PAD_CARTE
                for t in textes]

    # `taille_forcee` (taille commune calculee sur plusieurs colonnes) prime ; sinon
    # la colonne s'auto-ajuste (plus grande taille entre `small` et TAILLE_MIN_CARTE
    # telle que la pile tienne). Le forcage sert a accorder les slides 4 et 5.
    taille = taille_forcee if taille_forcee is not None else _taille_colonne(items, w)
    card_hs = hauteurs(taille)
    total = sum(card_hs)
    gap = max(GAP_MIN, min(GAP_MAX, (band - total) / (n - 1))) if n > 1 else 0.0
    total_avec_gaps = total + max(0, n - 1) * gap

    # Filet de securite : si meme a TAILLE_MIN_CARTE le total deborde encore
    # (texte extreme), on comprime proportionnellement les hauteurs de carte
    # pour GARANTIR qu'aucune forme ne deborde de la slide (invariant verifie
    # par D.verifier_geometrie). `ql_max` (deduit de la hauteur finale) fait
    # alors tronquer, dans le rendu, le texte de la carte concernee (ellipse) —
    # au pire un texte coupe, jamais une forme hors cadre.
    facteur = min(1.0, (band - 1e-6) / total_avec_gaps) if total_avec_gaps > 0 else 1.0
    card_hs = [h * facteur for h in card_hs]
    gap *= facteur

    lh = LH_QUESTION * (taille / D.TYPE["small"])
    y = top
    for it, card_h in zip(items, card_hs):
        ql_max = max(1, int((card_h - PAD_CARTE - _CARTE_H_FIXE) / lh + 1e-6))
        D.add_card(slide, x, y, w, card_h, accent)
        rendu(slide, x, y, w, card_h, it, taille, ql_max)
        y += card_h + gap


def slide_points(prs, layouts, bloc):
    slide = titre_slide(prs, layouts, f"{bloc['nom']} — Points d'attention")
    # On arrete les colonnes au bord droit "sur" (BORD_DROIT) pour degager le badge
    # n° de slide du template OCTO, que la carte du bas viendrait sinon toucher.
    colw = (BORD_DROIT - MARGE_X - 0.5) / 2
    xg, xd = MARGE_X, MARGE_X + colw + 0.5
    pad = 0.24                              # marge interne carte (apres le liseré)
    # Taille de carte COMMUNE aux slides 4 et 5 (meme calcul dans slide_points_forts).
    taille_cartes = _taille_cartes_bloc(bloc, colw)

    def rendu_dispersion(slide, x, y, w, h, q, taille, ql_max):
        tx = x + pad
        tw = w - pad - 0.18
        # Bloc {question + contexte + barre} centre dans la carte ; chaque element
        # suit le precedent (pas de vide au milieu, carte non etiree). `taille`
        # est choisie par _cartes_colonne pour que la question la plus longue de
        # la colonne tienne sans deborder (D.ajuster_police).
        texte, ql = _texte_et_lignes(q.get("texte", ""), tw, taille, ql_max)
        top0 = y + (h - _bloc_carte_h(ql, taille)) / 2
        qh = ql * LH_QUESTION * (taille / D.TYPE["small"])
        D.add_text(slide, tx, top0, tw, qh,
                   [(texte, {"size": taille, "bold": True,
                             "line_spacing": 0.96})])
        D.add_text(slide, tx, top0 + qh + 0.04, tw, 0.17,
                   [(_contexte_joli(q.get("contexte", "")),
                     {"size": D.TYPE["tiny"], "color": D.MUTED})])
        # Barre d'amplitude min..max sur l'echelle 0..3 + repere de moyenne.
        ry = top0 + qh + 0.31
        rw = w - pad - 1.55
        mn = q.get("min") if q.get("min") is not None else 0
        mx = q.get("max") if q.get("max") is not None else 3
        _widget_amplitude(slide, tx, ry, rw, mn, mx, q.get("moyenne"))
        # La plage min–max est deja montree par la barre ; on n'affiche que la
        # metrique de classement, nommee en clair (et non l'abreviation "é-t").
        _valeur_cote_barre(slide, tx + rw + 0.14, ry, w - pad - rw - 0.20,
                           [(fmt(q.get("ecartType")),
                             {"size": D.TYPE["h3"], "bold": True, "color": D.INK,
                              "align": PP_ALIGN.RIGHT}),
                            ("écart-type", {"size": D.TYPE["tiny"], "color": D.MUTED,
                                            "align": PP_ALIGN.RIGHT})])

    def rendu_faible(slide, x, y, w, h, q, taille, ql_max):
        tx = x + pad
        tw = w - pad - 0.18
        texte, ql = _texte_et_lignes(q.get("texte", ""), tw, taille, ql_max)
        top0 = y + (h - _bloc_carte_h(ql, taille, _CARTE_H_FIXE_SCORES)) / 2   # carte sans "moy." : centrer sur le contenu reel
        qh = ql * LH_QUESTION * (taille / D.TYPE["small"])
        D.add_text(slide, tx, top0, tw, qh,
                   [(texte, {"size": taille, "bold": True,
                             "line_spacing": 0.96})])
        D.add_text(slide, tx, top0 + qh + 0.04, tw, 0.17,
                   [(_contexte_joli(q.get("contexte", "")),
                     {"size": D.TYPE["tiny"], "color": D.MUTED})])
        ry = top0 + qh + 0.31
        rw = w - pad - 1.35
        moy = q.get("moyenne")
        D.add_hbar(slide, tx, ry, rw, 0.13, (moy / 3.0) if moy is not None else 0, D.INK)
        _valeur_cote_barre(slide, tx + rw + 0.14, ry, w - pad - rw - 0.20,
                           [(fmt(moy), {"size": D.TYPE["h3"], "bold": True, "color": D.INK,
                                        "align": PP_ALIGN.RIGHT}),
                            ("sur 3", {"size": D.TYPE["tiny"], "color": D.MUTED,
                                       "align": PP_ALIGN.RIGHT})])

    _entete_colonne(slide, xg, colw, "▼", "Plus forts désaccords",
                    "Forte dispersion des réponses — sujets à clarifier")
    _cartes_colonne(slide, xg, colw, bloc.get("dispersion", []), D.INK, rendu_dispersion, taille_cartes)
    _entete_colonne(slide, xd, colw, "▼", "Scores les plus faibles",
                    "Maturité la plus basse — leviers de progrès prioritaires")
    _cartes_colonne(slide, xd, colw, bloc.get("faibles", []), D.INK, rendu_faible, taille_cartes)


# ----------------------------------------------------------------------------
# Slide "Points forts" : pendant positif de la slide "Points d'attention" —
# scores les plus hauts et meilleurs accords (dispersion la plus faible).
# ----------------------------------------------------------------------------
def slide_points_forts(prs, layouts, bloc):
    slide = titre_slide(prs, layouts, f"{bloc['nom']} — Points forts")
    colw = (BORD_DROIT - MARGE_X - 0.5) / 2
    xg, xd = MARGE_X, MARGE_X + colw + 0.5
    pad = 0.24
    # Taille de carte COMMUNE aux slides 4 et 5 (meme calcul dans slide_points) :
    # les deux slides affichent ainsi la meme taille de texte.
    taille_cartes = _taille_cartes_bloc(bloc, colw)

    def rendu_haut(slide, x, y, w, h, q, taille, ql_max):
        tx = x + pad
        tw = w - pad - 0.18
        texte, ql = _texte_et_lignes(q.get("texte", ""), tw, taille, ql_max)
        top0 = y + (h - _bloc_carte_h(ql, taille, _CARTE_H_FIXE_SCORES)) / 2   # carte sans "moy." : centrer sur le contenu reel
        qh = ql * LH_QUESTION * (taille / D.TYPE["small"])
        D.add_text(slide, tx, top0, tw, qh,
                   [(texte, {"size": taille, "bold": True, "line_spacing": 0.96})])
        D.add_text(slide, tx, top0 + qh + 0.04, tw, 0.17,
                   [(_contexte_joli(q.get("contexte", "")),
                     {"size": D.TYPE["tiny"], "color": D.MUTED})])
        ry = top0 + qh + 0.31
        rw = w - pad - 1.35
        moy = q.get("moyenne")
        D.add_hbar(slide, tx, ry, rw, 0.13, (moy / 3.0) if moy is not None else 0, D.INK)
        _valeur_cote_barre(slide, tx + rw + 0.14, ry, w - pad - rw - 0.20,
                           [(fmt(moy), {"size": D.TYPE["h3"], "bold": True, "color": D.INK,
                                        "align": PP_ALIGN.RIGHT}),
                            ("sur 3", {"size": D.TYPE["tiny"], "color": D.MUTED,
                                       "align": PP_ALIGN.RIGHT})])

    def rendu_accord(slide, x, y, w, h, q, taille, ql_max):
        tx = x + pad
        tw = w - pad - 0.18
        texte, ql = _texte_et_lignes(q.get("texte", ""), tw, taille, ql_max)
        top0 = y + (h - _bloc_carte_h(ql, taille)) / 2
        qh = ql * LH_QUESTION * (taille / D.TYPE["small"])
        D.add_text(slide, tx, top0, tw, qh,
                   [(texte, {"size": taille, "bold": True, "line_spacing": 0.96})])
        D.add_text(slide, tx, top0 + qh + 0.04, tw, 0.17,
                   [(_contexte_joli(q.get("contexte", "")),
                     {"size": D.TYPE["tiny"], "color": D.MUTED})])
        ry = top0 + qh + 0.31
        rw = w - pad - 1.55
        mn = q.get("min") if q.get("min") is not None else 0
        mx = q.get("max") if q.get("max") is not None else 3
        _widget_amplitude(slide, tx, ry, rw, mn, mx, q.get("moyenne"))
        _valeur_cote_barre(slide, tx + rw + 0.14, ry, w - pad - rw - 0.20,
                           [(fmt(q.get("ecartType")),
                             {"size": D.TYPE["h3"], "bold": True, "color": D.INK,
                              "align": PP_ALIGN.RIGHT}),
                            ("écart-type", {"size": D.TYPE["tiny"], "color": D.MUTED,
                                            "align": PP_ALIGN.RIGHT})])

    _entete_colonne(slide, xg, colw, "▲", "Scores les plus hauts",
                    "Maturité la plus haute — points d'appui à valoriser")
    _cartes_colonne(slide, xg, colw, bloc.get("hauts", []), D.INK, rendu_haut, taille_cartes)
    _entete_colonne(slide, xd, colw, "▲", "Meilleurs accords",
                    "Dispersion la plus faible des réponses — consensus fort")
    accords = bloc.get("accords", [])
    if accords:
        _cartes_colonne(slide, xd, colw, accords, D.INK, rendu_accord, taille_cartes)
    else:
        # Un accord n'a de sens qu'avec >= 2 reponses (ex. equipe a 1 seul
        # repondant) : etat vide explicite plutot qu'une colonne silencieuse.
        D.add_text(slide, xd + 0.05, CONTENU_TOP + 0.66, colw - 0.10, 0.4,
                   [("Pas assez de réponses pour mesurer un accord.",
                     {"size": D.TYPE["small"], "italic": True, "color": D.MUTED})])


# ----------------------------------------------------------------------------
def construire(data, template_path, out_path):
    prs = Presentation(template_path)
    # Charte du modele : on prend la couleur primaire du theme (dk1) comme accent
    # de marque ; repli sur le bleu de la palette si le theme est illisible.
    global ACCENT, CYAN, FOND_PANNEAU
    # Aligne tout le deck sur la charte du template : police de marque (Outfit) +
    # neutres navy/slate + accent cyan, lus dans le theme (= charte OCTO). Sans ca
    # le texte heritait d'Arial et les neutres etaient des gris generiques codes en
    # dur. Detecte, pas code en dur : s'adapte a un autre template fourni.
    marque = D.appliquer_theme(prs)
    ACCENT = marque.get("navy") or D.PALETTE[0]   # accent principal = navy (jauge)
    CYAN = marque.get("cyan") or ACCENT           # accent secondaire = cyan
    FOND_PANNEAU = D.TRACK                          # fond d'encart = slate 100 du theme
    layouts = prs.slide_masters[0].slide_layouts
    nb_template = len(prs.slides)

    cv = data.get("couverture")
    if cv:
        slide = prs.slides.add_slide(_trouver_layout(layouts, COUV_PATTERNS, LAYOUT_COUVERTURE))
        ph = {p.placeholder_format.idx: p for p in slide.placeholders}
        if 0 in ph: ph[0].text_frame.text = cv.get("titre", "Restitution")
        if 1 in ph: ph[1].text_frame.text = cv.get("sousTitre", "")
        if 2 in ph: ph[2].text_frame.text = "OCTO Technology"
        if 3 in ph: ph[3].text_frame.text = cv.get("date", "")

    for bloc in data.get("blocs", []):
        slide_vue_ensemble(prs, layouts, bloc)
        slide_radar(prs, layouts, bloc)
        slide_progression(prs, layouts, bloc)
        slide_points_forts(prs, layouts, bloc)
        slide_points(prs, layouts, bloc)

    # Retire les slides d'exemple du template (garde masters/layouts OCTO).
    xml_slides = prs.slides._sldIdLst
    for slide_xml in list(xml_slides)[:nb_template]:
        xml_slides.remove(slide_xml)

    problemes = D.verifier_geometrie(prs)
    if problemes:
        sys.stderr.write("ATTENTION geometrie :\n" + "\n".join(problemes) + "\n")

    prs.save(out_path)
    return prs, problemes


if __name__ == "__main__":
    if len(sys.argv) < 3:
        sys.stderr.write(
            "usage: python export-restitution-ppt.py "
            "<donnees.json> <sortie.pptx> [modele.pptx]\n")
        sys.exit(2)
    with open(sys.argv[1], encoding="utf-8") as f:
        data = json.load(f)
    # Template du modele : 3e argument, sinon $TEMPLATE_PPTX, sinon le template OCTO.
    # Le serveur n'en passe que 2 (data, out) => repli sur TEMPLATE (compat preservee).
    template = (sys.argv[3] if len(sys.argv) > 3 else None) or os.environ.get("TEMPLATE_PPTX") or TEMPLATE
    prs, problemes = construire(data, template, sys.argv[2])
    print(sys.argv[2], "-", len(prs.slides), "slides",
          "- geometrie OK" if not problemes else f"- {len(problemes)} pb geometrie")
