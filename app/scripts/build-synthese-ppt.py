import copy
import datetime
from pptx import Presentation
from pptx.util import Inches, Pt
from PIL import Image

TEMPLATE_PATH = "template ppt/template.pptx"
OUTPUT_PATH = "template ppt/Synthese-MVP-Increment2.pptx"
CAPTURES_DIR = "cadrage/captures"

LAYOUT_COUVERTURE = 8
LAYOUT_TITRE_CONTENU = 3

prs = Presentation(TEMPLATE_PATH)
layouts = prs.slide_masters[0].slide_layouts
NB_SLIDES_TEMPLATE_A_RETIRER = len(prs.slides)


def set_text(placeholder, text, size=None, bold=None):
    placeholder.text_frame.text = text
    if size or bold is not None:
        for paragraph in placeholder.text_frame.paragraphs:
            for run in paragraph.runs:
                if size:
                    run.font.size = Pt(size)
                if bold is not None:
                    run.font.bold = bold


def slide_couverture(titre, sous_titre):
    slide = prs.slides.add_slide(layouts[LAYOUT_COUVERTURE])
    phs = {ph.placeholder_format.idx: ph for ph in slide.placeholders}
    set_text(phs[0], titre)
    set_text(phs[1], sous_titre)
    set_text(phs[2], "OCTO Technology")
    set_text(phs[3], datetime.date.today().strftime("%d.%m.%y"))
    return slide


def _remplir_colonne(text_frame, entete, puces, taille_pt=13):
    text_frame.word_wrap = True
    text_frame.text = entete
    text_frame.paragraphs[0].runs[0].font.bold = True
    text_frame.paragraphs[0].runs[0].font.size = Pt(14)
    for puce in puces:
        p = text_frame.add_paragraph()
        p.text = puce
        p.runs[0].font.size = Pt(taille_pt)


def slide_executive_summary(titre, realise, reste_a_faire):
    slide = prs.slides.add_slide(layouts[LAYOUT_TITRE_CONTENU])
    phs = {ph.placeholder_format.idx: ph for ph in slide.placeholders}
    set_text(phs[0], titre)

    colonne_gauche = phs[1]
    colonne_gauche.left = Inches(0.63)
    colonne_gauche.top = Inches(0.95)
    colonne_gauche.width = Inches(4.1)
    colonne_gauche.height = Inches(4.2)
    _remplir_colonne(colonne_gauche.text_frame, "Réalisé (Increment 1 & 2)", realise)

    # Duplique le placeholder (plutot qu'une textbox libre) pour conserver les
    # puces du style de liste du theme, qu'une simple zone de texte n'a pas.
    nouvel_element = copy.deepcopy(colonne_gauche._element)
    slide.shapes._spTree.append(nouvel_element)
    colonne_droite = next(s for s in slide.shapes if s._element is nouvel_element)
    colonne_droite.left = Inches(5.2)
    colonne_droite.top = Inches(0.95)
    colonne_droite.width = Inches(4.1)
    colonne_droite.height = Inches(4.2)
    _remplir_colonne(colonne_droite.text_frame, "Reste à faire", reste_a_faire)

    return slide


def slide_capture(titre, puces, image_path):
    slide = prs.slides.add_slide(layouts[LAYOUT_TITRE_CONTENU])
    phs = {ph.placeholder_format.idx: ph for ph in slide.placeholders}
    set_text(phs[0], titre)

    legende = phs[1]
    legende.left = Inches(0.63)
    legende.top = Inches(0.95)
    legende.width = Inches(2.5)
    legende.height = Inches(4.2)
    tf = legende.text_frame
    tf.word_wrap = True
    tf.text = puces[0]
    for puce in puces[1:]:
        p = tf.add_paragraph()
        p.text = puce
    for paragraph in tf.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(13)

    zone_left = Inches(3.45)
    zone_top = Inches(0.95)
    zone_max_width = Inches(5.9)
    zone_max_height = Inches(4.2)

    with Image.open(image_path) as im:
        aspect = im.width / im.height

    width = zone_max_width
    height = int(width / aspect)
    if height > zone_max_height:
        height = zone_max_height
        width = int(height * aspect)

    top = int(zone_top + (zone_max_height - height) / 2)
    left = int(zone_left + (zone_max_width - width) / 2)
    slide.shapes.add_picture(image_path, left, top, width=width, height=height)
    return slide


# 1. Couverture
slide_couverture(
    "Synthèse du réalisé",
    "Questionnaire de maturité agile/produit — MVP, Increment 1 & 2",
)

# 2. Executive summary
slide_executive_summary(
    "Executive Summary",
    [
        "Increment 1 — Collecte : import/ré-import du référentiel Excel, création de session, identification du répondant, parcours de réponse complet, soumission verrouillée.",
        "Increment 2 — Résultats bruts : écran animateur agrégé/anonyme avec radar par objectif, drill-down nominatif par question, filtre avec/sans manager.",
        "Correction orthographique automatique des textes importés (vocabulaire agile préservé).",
        "Hiérarchie pilier → objectif → question reconstruite directement depuis le fichier Excel source.",
    ],
    [
        "Increment 3 — Restitution : pré-analyses (minimum, maximum, dispersion, taux de réponse).",
        "Export du support de restitution au format slide (PPT).",
        "Comparaison multi-sessions : radar superposé, régressions/progressions, gestion du turnover.",
        "Epic 7 (hors MVP) — Consolidation multi-équipes : vue dédiée Sponsor/RH/Direction.",
    ],
)

# 4. Capture : espace animateur
slide_capture(
    "Espace animateur — Import & session",
    [
        "Import du fichier Excel (1 onglet = 1 pilier).",
        "Un nouvel import remplace intégralement le référentiel.",
        "Création d'une session avec dates d'ouverture/fermeture de la saisie.",
    ],
    f"{CAPTURES_DIR}/01-admin.png",
)

# 5. Capture : parcours répondant
slide_capture(
    "Parcours répondant — Questionnaire",
    [
        "Navigation par pilier, progression visible.",
        "Réponse par \"Choix 0-3\" sans terminologie de notation visible.",
        "Détail du choix affiché dans un encart dédié.",
    ],
    f"{CAPTURES_DIR}/02-repondre.png",
)

# 6. Capture : résultats - radar
slide_capture(
    "Résultats animateur — Radar par objectif",
    [
        "Un axe par objectif (sous-catégorie), coloré selon son pilier.",
        "Filtre avec/sans manager.",
        "Calcul agrégé en temps réel.",
    ],
    f"{CAPTURES_DIR}/03-resultats-radar.png",
)

# 7. Capture : résultats - détail nominatif
slide_capture(
    "Résultats animateur — Détail par pilier",
    [
        "Accordéon pilier → objectif → question.",
        "Drill-down nominatif à la demande, par question.",
        "Aucun seuil minimal d'anonymat.",
    ],
    f"{CAPTURES_DIR}/04-resultats-detail.png",
)

# Retire les slides d'exemple du template (mission HERMES) : on ne garde
# que les masters/layouts (branding OCTO) et les slides nouvellement ajoutees.
xml_slides = prs.slides._sldIdLst
slides_a_retirer = list(xml_slides)[:NB_SLIDES_TEMPLATE_A_RETIRER]
for slide_xml in slides_a_retirer:
    xml_slides.remove(slide_xml)

prs.save(OUTPUT_PATH)
print("PPT genere :", OUTPUT_PATH, "-", len(prs.slides), "slides")
