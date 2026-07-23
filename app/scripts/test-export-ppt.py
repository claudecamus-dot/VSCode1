"""Test du generateur de restitution PPT (US6.4).

Verifie, sans LibreOffice, que le deck genere :
  - a le bon nombre de slides (1 couverture + 5 par bloc : vue d'ensemble, radar,
    progression, points forts, points d'attention) ;
  - ne contient AUCUNE forme hors cadre (pptx_deck.verifier_geometrie), y
    compris avec une question anormalement longue (auto-ajustement de police) ;
  - se construit aussi bien avec qu'avec radar large / sans comparaison /
    valeurs manquantes (robustesse).

Usage : python test-export-ppt.py
"""
import os
import sys
import struct
import zlib
import tempfile

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
# Le generateur a un nom de fichier non importable (tirets) -> import par chemin.
import importlib.util
_spec = importlib.util.spec_from_file_location(
    "gen", os.path.join(os.path.dirname(__file__), "export-restitution-ppt.py"))
gen = importlib.util.module_from_spec(_spec)
_spec.loader.exec_module(gen)

echecs = 0


def check(cond, msg):
    global echecs
    if cond:
        print(f"  ok   {msg}")
    else:
        echecs += 1
        print(f"  FAIL {msg}")


def png_factice(path, w, h):
    """Ecrit un PNG uni minimal de dimensions w x h (sans PIL)."""
    def chunk(typ, data):
        c = typ + data
        return struct.pack(">I", len(data)) + c + struct.pack(">I", zlib.crc32(c) & 0xffffffff)
    sig = b"\x89PNG\r\n\x1a\n"
    ihdr = struct.pack(">IIBBBBB", w, h, 8, 2, 0, 0, 0)
    raw = b"".join(b"\x00" + b"\x9c\xc0\xd6" * w for _ in range(h))
    idat = zlib.compress(raw)
    with open(path, "wb") as fp:
        fp.write(sig + chunk(b"IHDR", ihdr) + chunk(b"IDAT", idat) + chunk(b"IEND", b""))


def piliers(vals):
    noms = ["Equipe Produit", "Excellence Technique", "Culture De L'Entreprise Agile", "Agilite A L'Echelle"]
    return [{"nom": n, "moyenne": v} for n, v in zip(noms, vals)]


# Deux libelles avec un complement entre parentheses (comme le referentiel
# Excel reel, ex. "Ressources humaines (formations, coaching agile, talent,
# ...)") : doit s'afficher SANS la parenthese, partout ou le nom est affiche
# (labels d'axe du radar ET legende) — regression deja signalee une fois.
_NOMS_OBJECTIFS = [
    ["Vision produit", "Découpage en user stories",
     "Ressources humaines (formations, coaching agile, talent, ...)"],
    ["Qualité et tests automatisés", "Intégration continue", "Architecture évolutive"],
    ["Sponsorship et posture managériale", "Droit à l'erreur", "Collaboration transverse"],
    ["Synchronisation inter-équipes",
     "Fonctionnement agile à l'échelle (gestion des dépendances, méthodes, outils, ...)",
     "Gouvernance à l'échelle"],
]


def objectifs(vals, avec_precedent=True):
    """3 sous-categories par pilier (radar a 12 axes, comme un referentiel
    reel) — dont un libelle anormalement long (stress-test du repli de police/
    wrap du radar vectoriel) et, si `avec_precedent`, un axe SANS comparaison
    (doit retomber sur sa valeur courante, pas s'effondrer a 0)."""
    axes = []
    for pilier_i, (noms, v) in enumerate(zip(_NOMS_OBJECTIFS, vals)):
        for j, nom in enumerate(noms):
            sans_comp = avec_precedent and pilier_i == 1 and j == 1
            axes.append({
                "nom": nom,
                "moyenne": max(0.0, min(3.0, v + (j - 1) * 0.3)),
                "precedent": None if (not avec_precedent or sans_comp) else max(0.0, v - 0.6),
                "pilierIndex": pilier_i,
            })
    return axes


def bloc_equipe(nom, vals, avec_comp=True, radar=None):
    b = {
        "type": "equipe", "nom": nom, "departement": "DSI Paiements", "effectif": 5,
        "piliers": piliers(vals),
        "objectifs": objectifs(vals, avec_precedent=avec_comp),
        "dispersion": [
            {"texte": "Est-ce que le recrutement est adapte aux besoins de l'equipe ?",
             "ecartType": 1.0, "min": 1, "max": 3, "moyenne": 2.0, "contexte": "RH"},
            {"texte": "Comment les utilisateurs sont-ils consultes regulierement et tot ?",
             "ecartType": 0.9, "min": 0, "max": 3, "moyenne": 1.5, "contexte": "Produit"},
            {"texte": "Comment est mise en oeuvre l'amelioration continue dans l'equipe ?",
             "ecartType": 0.8, "min": 1, "max": 3, "moyenne": 2.2, "contexte": "Culture"},
        ],
        "faibles": [
            {"texte": "Les besoins utilisateurs sont-ils compris et reformules en user stories ?",
             "moyenne": 1.4, "contexte": "Produit"},
            {"texte": "Les taches sont-elles decrites, decoupees et ordonnees dans un backlog ?",
             "moyenne": 1.6, "contexte": "Produit"},
            # Question anormalement longue : stress-test de l'auto-ajustement de
            # police (D.ajuster_police) — ne doit ni deborder ni etre tronquee.
            {"texte": ("L'equipe a-t-elle les moyens d'organiser son delivery pour livrer "
                       "regulierement de la valeur aux utilisateurs finaux, en tenant compte "
                       "des contraintes de securite, de conformite et d'architecture technique "
                       "imposees par le systeme d'information de l'entreprise ?"),
             "moyenne": 1.6, "contexte": "Technique"},
        ],
        "hauts": [
            {"texte": "Comment la strategie de release apporte-t-elle de la valeur aux clients ?",
             "moyenne": 2.8, "contexte": "Produit"},
            {"texte": "Les retrospectives ont-elles un impact reel sur le fonctionnement ?",
             "moyenne": 2.6, "contexte": "Culture"},
            {"texte": "L'equipe priorise-t-elle son backlog selon la valeur metier ?",
             "moyenne": 2.5, "contexte": "Produit"},
        ],
        "accords": [
            {"texte": "L'equipe sait-elle pourquoi elle livre ce qu'elle livre ?",
             "ecartType": 0.2, "min": 2, "max": 3, "contexte": "Vision"},
            {"texte": "Le role de chacun dans l'equipe est-il clair ?",
             "ecartType": 0.3, "min": 2, "max": 3, "contexte": "Culture"},
        ],
        "commentaire": "Nette progression depuis janvier sur le fonctionnement produit.\nVigilance : ingenierie/testing en retrait.",
    }
    if radar:
        b["radarImage"] = radar
    if avec_comp:
        b["comparaison"] = {
            "disponible": True, "precedenteDate": "15/01/2026",
            "piliers": [{"nom": p["nom"], "courant": p["moyenne"],
                         "precedent": max(0, p["moyenne"] - 0.8),
                         "delta": 0.8} for p in piliers(vals)],
        }
    else:
        b["comparaison"] = {"disponible": False}
    return b


def main():
    tmp = tempfile.mkdtemp(prefix="test-ppt-")
    # radarImage n'est plus utilise par le generateur (radar desormais vectoriel,
    # dessine depuis objectifs/piliers — voir _dessiner_radar) : ces PNG factices
    # ne servent plus qu'a verifier qu'un champ radarImage encore envoye par le
    # serveur (avant nettoyage cote server.js) reste sans effet, sans casser.
    radar_carre = os.path.join(tmp, "radar.png")
    radar_large = os.path.join(tmp, "radar-wide.png")
    png_factice(radar_carre, 520, 520)
    png_factice(radar_large, 900, 360)  # tres large : ne doit pas deborder

    dep = bloc_equipe("DSI Paiements", [2.0, 1.8, 1.6, 1.9], avec_comp=False, radar=radar_carre)
    dep["type"] = "departement"; dep["nbEquipes"] = 2; dep["departement"] = "DSI Paiements"
    data = {
        "couverture": {"titre": "Restitution — Maturite agile/produit",
                       "sousTitre": "Departement DSI Paiements", "date": "24/06/2026"},
        "blocs": [
            dep,
            bloc_equipe("Squad Paiement", [2.7, 2.0, 1.4, 2.3], radar=radar_large),
            bloc_equipe("Squad Virement", [1.5, 1.2, 1.0, 1.4], avec_comp=True, radar=radar_carre),
        ],
    }

    out = os.path.join(tmp, "deck.pptx")
    prs, problemes = gen.construire(data, gen.TEMPLATE, out)

    print("Structure :")
    check(len(prs.slides) == 1 + 3 * 5, f"16 slides (couverture + 3x5 : vue/radar/progression/forts/attention) — recu {len(prs.slides)}")
    check(os.path.exists(out) and os.path.getsize(out) > 0, "fichier .pptx ecrit")

    print("Geometrie (aucune forme hors cadre) :")
    if problemes:
        for p in problemes:
            print("   -", p)
    check(not problemes, f"toutes les formes dans le cadre — {len(problemes)} probleme(s)")

    print("Texte dans sa boite (verifier_debordements_texte, porte de VSCode2) :")
    # 1) Contrat de la fonction (test dur) : une boite trop basse pour son texte
    #    est flaguee, une boite dimensionnee correctement ne l'est pas.
    from pptx import Presentation
    prs_unit = Presentation()
    slide_unit = prs_unit.slides.add_slide(prs_unit.slide_layouts[6])
    long_txt = ("Un texte volontairement long qui doit se replier sur de "
                "nombreuses lignes dans une boite bien trop basse pour lui.")
    gen.D.add_text(slide_unit, 0.5, 0.5, 2.0, 0.2, [(long_txt, {"size": 12})])
    gen.D.add_text(slide_unit, 0.5, 2.0, 6.0, 1.0, [("Texte court.", {"size": 12})])
    constats_unit = gen.D.verifier_debordements_texte(prs_unit)
    check(len(constats_unit) == 1 and "slide 1" in constats_unit[0],
          f"le debordement volontaire est flague, la boite saine non — {constats_unit}")
    # 2) Filet sur le deck complet : AVERTISSEMENT seulement pour l'instant.
    #    L'estimateur est pessimiste par contrat ; le rendu actuel a ete valide a
    #    l'oeil par l'utilisateur, donc un constat ici = candidat a trier par
    #    rendu reel (skill deck-design-review), pas un echec automatique.
    debords = gen.D.verifier_debordements_texte(prs)
    if debords:
        print(f"   (avertissement : {len(debords)} boite(s) limite(s) a trier par rendu reel)")
        for p in debords:
            print("   -", p)
    else:
        print("   aucun constat")

    print("Libelles pilier/objectif : jamais de parenthese affichee (verrou anti-regression) :")
    tout_le_texte = []
    for slide in prs.slides:
        for shp in slide.shapes:
            if shp.has_text_frame:
                tout_le_texte.append(shp.text_frame.text)
    texte_complet = "\n".join(tout_le_texte)
    check("(formations" not in texte_complet,
          "pas de complement parenthese residuel (ex. 'Ressources humaines (formations...')")
    check("(gestion des" not in texte_complet,
          "pas de complement parenthese residuel (ex. 'Fonctionnement agile ... (gestion des...')")

    print("Robustesse — bloc sans comparaison et valeurs manquantes :")
    data2 = {"couverture": None, "blocs": [{
        "type": "equipe", "nom": "Equipe Vide", "departement": "", "effectif": 0,
        "piliers": [{"nom": "P1", "moyenne": None}, {"nom": "P2", "moyenne": 2.0}],
        "objectifs": [], "dispersion": [], "faibles": [], "hauts": [], "accords": [],
        "commentaire": "", "comparaison": {"disponible": False},
    }]}
    out2 = os.path.join(tmp, "deck2.pptx")
    prs2, pb2 = gen.construire(data2, gen.TEMPLATE, out2)
    check(len(prs2.slides) == 5, f"5 slides (pas de couverture) — recu {len(prs2.slides)}")
    check(not pb2, f"geometrie OK meme avec donnees partielles — {len(pb2)} probleme(s)")

    print("Robustesse — radar vectoriel avec trop peu d'axes (< 3, illisible) :")
    data3 = {"couverture": None, "blocs": [{
        "type": "equipe", "nom": "Equipe Deux Axes", "departement": "", "effectif": 3,
        "piliers": [{"nom": "P1", "moyenne": 2.0}, {"nom": "P2", "moyenne": 1.0}],
        "objectifs": [{"nom": "Obj A", "moyenne": 2.0, "precedent": None, "pilierIndex": 0},
                      {"nom": "Obj B", "moyenne": 1.0, "precedent": None, "pilierIndex": 1}],
        "dispersion": [], "faibles": [], "hauts": [], "accords": [],
        "commentaire": "Radar non affiche (2 axes seulement).", "comparaison": {"disponible": False},
    }]}
    out3 = os.path.join(tmp, "deck3.pptx")
    prs3, pb3 = gen.construire(data3, gen.TEMPLATE, out3)
    check(not pb3, f"geometrie OK avec seulement 2 axes (radar non dessine) — {len(pb3)} probleme(s)")

    print("Robustesse — comparaison disponible + < 3 objectifs (pas de radar), et max=0 :")
    # Regression : radar absent (axes<3) => l'ancien repli px=7.4 coincait le
    # panneau evolution a droite (pw=2.05) et produisait des textboxes de largeur
    # NEGATIVE (name_w=-0.22), invisibles a l'ancien verifier_geometrie (qui ne
    # testait que le debordement des bords). Le test a 2 axes existant (data3)
    # force comparaison.disponible=false et ratait donc ce chemin. Ici on force la
    # comparaison + des cartes max=0 (barre d'amplitude qui doit rester une
    # pastille a 0, pas une barre pleine — bug `q.get("max",3) or 3`).
    data4 = {"couverture": None, "blocs": [{
        "type": "equipe", "nom": "Equipe Deux Axes Comparee", "departement": "", "effectif": 4,
        "piliers": [{"nom": "Pilier Un", "moyenne": 2.0}, {"nom": "Pilier Deux", "moyenne": 0.0}],
        "objectifs": [{"nom": "Obj A", "moyenne": 2.0, "precedent": 1.5, "pilierIndex": 0},
                      {"nom": "Obj B", "moyenne": 0.0, "precedent": 0.5, "pilierIndex": 1}],
        "dispersion": [{"texte": "Tous au niveau zero sur ce sujet peu mature ?",
                        "ecartType": 0.0, "min": 0, "max": 0, "moyenne": 0.0, "contexte": "Culture"}],
        "faibles": [], "hauts": [],
        "accords": [{"texte": "Consensus total, au niveau zero ?",
                     "ecartType": 0.0, "min": 0, "max": 0, "contexte": "Vision"}],
        "commentaire": "Deux axes seulement, mais une session precedente existe.",
        "comparaison": {"disponible": True, "precedenteDate": "01/01/2026",
                        "piliers": [{"nom": "Pilier Un", "courant": 2.0, "precedent": 1.2, "delta": 0.8},
                                    {"nom": "Pilier Deux", "courant": 0.0, "precedent": 0.3, "delta": -0.3}]},
    }]}
    out4 = os.path.join(tmp, "deck4.pptx")
    prs4, pb4 = gen.construire(data4, gen.TEMPLATE, out4)
    if pb4:
        for p in pb4:
            print("   -", p)
    check(not pb4, f"aucune forme degeneree (largeur negative) avec comparaison + 2 axes + max=0 — {len(pb4)} probleme(s)")

    print("\nTOUS LES TESTS PASSENT" if echecs == 0 else f"\n{echecs} TEST(S) EN ECHEC")
    sys.exit(0 if echecs == 0 else 1)


if __name__ == "__main__":
    main()
